# CouSolve Report 양식(7단계 폼)에 내용만 채워 넣기 위한 도구.
# 폼 자체 — 상단 단계 칩 · 둥근 컨테이너 · 초록 라벨 · 표 틀 — 은 손대지 않는다.
# 글꼴과 색은 양식이 쓰던 것을 그대로 따른다 (에스코어 드림 / #0D4B45).
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

GREEN = RGBColor(0x0D, 0x4B, 0x45)   # 양식 Theme 바 색
GREEN2 = RGBColor(0x10, 0x5C, 0x4D)  # 양식 구역 라벨 색
INK = RGBColor(0x00, 0x00, 0x00)
INK2 = RGBColor(0x44, 0x4F, 0x55)
INK3 = RGBColor(0x80, 0x8A, 0x90)
LINE = RGBColor(0xD5, 0xDB, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
FONT = "에스코어 드림 4 Regular"
FONT_M = "에스코어 드림 5 Medium"


def _run(para, text, size, bold=False, color=INK, italic=False, font=None, u=False):
    """u=True 는 추정치 표시다. 자료 전체에서 밑줄은 '실측 아님' 하나의 뜻만 갖는다."""
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.underline = bool(u)
    r.font.name = font or (FONT_M if bold else FONT)
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = parse_xml(f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            rPr.append(e)
        e.set("typeface", r.font.name)
    return r


def _chunks(para, chunks, size, color):
    """chunks: 문자열이거나 [(글자, {옵션}), ...]"""
    if isinstance(chunks, str):
        chunks = [(chunks, {})]
    for t, o in chunks:
        _run(para, t, o.get("size", size), o.get("bold", False),
             o.get("color", color), o.get("italic", False), o.get("font"), o.get("u", False))


def text(slide, x, y, w, h, chunks, size=10, align=PP_ALIGN.LEFT, color=INK,
         space=0, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if space:
        p.line_spacing = space
    _chunks(p, chunks, size, color)
    return box


def lines(slide, x, y, w, items, size=10, gap=0.30, color=INK, bullet="·  "):
    """한 줄짜리 항목을 일정 간격으로 쌓는다. items 는 문자열 또는 chunk 리스트."""
    for i, it in enumerate(items):
        c = it if not isinstance(it, str) else [(bullet + it, {})]
        text(slide, x, y + i * gap, w, gap, c, size=size, color=color)
    return y + len(items) * gap


def cell(c, chunks, size=10, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
         color=INK, margin=0.07, space=0):
    """이미 있는 표 칸에 글자만 채운다 — 칸 배경·테두리는 양식 것을 그대로 둔다."""
    tf = c.text_frame
    tf.word_wrap = True
    c.margin_left = c.margin_right = Inches(margin)
    c.margin_top = c.margin_bottom = Inches(0.02)
    c.vertical_anchor = anchor
    body = tf.paragraphs[0]._p.getparent()
    NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for p in list(body.findall(NS + "p"))[1:]:
        body.remove(p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    rows = chunks if isinstance(chunks, list) and chunks and isinstance(chunks[0], list) else [chunks]
    for i, row in enumerate(rows):
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.line_spacing = space
        _chunks(p, row, size, color)
    return c


def _border(c, edge, color, pt=0.75):
    tag = {"L": "a:lnL", "R": "a:lnR", "T": "a:lnT", "B": "a:lnB"}[edge]
    tcPr = c._tc.get_or_add_tcPr()
    for e in tcPr.findall(qn(tag)):
        tcPr.remove(e)
    ln = parse_xml(
        f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'w="{int(pt * 12700)}" cap="flat" cmpd="sng" algn="ctr">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'<a:prstDash val="solid"/></{tag}>')
    order = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
    idx = order.index(tag)
    after = None
    for t in order[idx + 1:]:
        after = tcPr.find(qn(t))
        if after is not None:
            break
    if after is not None:
        after.addprevious(ln)
    else:
        tcPr.insert(0, ln)


def table(slide, x, y, w, header, rows, col_w, row_h=0.28, head_h=0.28,
          size=9, head_size=9, aligns=None, colors=None, head_fill=GREEN):
    """양식에 없는 자리에 새 표를 세운다. 머리행은 Theme 초록."""
    n_r, n_c = len(rows) + 1, len(header)
    gf = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w),
                                Inches(head_h + row_h * len(rows)))
    tbl = gf.table
    tbl._tbl.find(qn("a:tblPr")).set("bandRow", "0")
    tbl._tbl.find(qn("a:tblPr")).set("firstRow", "1")
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r):
        tbl.rows[i].height = Inches(row_h)
    for j, h in enumerate(header):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = head_fill
        cell(c, [(h, {"bold": True, "color": WHITE})], head_size, PP_ALIGN.CENTER, margin=0.05)
        for e in "LRTB":
            _border(c, e, "0D4B45")
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            al = (aligns or [PP_ALIGN.LEFT] * n_c)[j]
            col = (colors or [INK] * n_c)[j]
            cell(c, v if not isinstance(v, str) else [(v, {})], size, al, margin=0.05, color=col)
            for e in "LRT":
                _border(c, e, "FFFFFF", 0.75)
            _border(c, "B", "D5DBDB", 0.75)
    return gf


def zone(slide, x, y, w, cols, size=9.5, gap=0.255, head_gap=0.30):
    """구역 라벨 아래에 '소제목 + 항목' 묶음을 열로 배치한다."""
    n = len(cols)
    cw = w / n
    bottom = y
    for i, (head, items) in enumerate(cols):
        cx = x + i * cw
        if head:
            text(slide, cx, y, cw - 0.2, 0.26,
                 [("▪ ", {"color": GREEN2, "bold": True}), (head, {"bold": True, "color": GREEN2})],
                 size=size + 0.5)
            yy = y + head_gap
        else:
            yy = y
        for it in items:
            c = it if not isinstance(it, str) else [("·  " + it, {})]
            text(slide, cx + (0.10 if head else 0), yy, cw - 0.24, gap, c, size=size)
            yy += gap
        bottom = max(bottom, yy)
    return bottom


def chart(slide, x, y, w, h, cats, vals, kind="col", num_fmt="#,##0",
          label_size=10, cat_size=9, color=GREEN, gap=110, headroom=1.18):
    """단일 계열 막대. 값 축·눈금선은 지우고 데이터 레이블로 직접 읽힌다.

    단위는 레이블에 붙이지 않고 캡션에 적는다 — 막대 폭에 눌려 줄바꿈된다.
    company_form/_devices.py 에 같은 함수가 있다. 두 덱은 서로 독립이라 색만 바꿔 따로 둔다.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK

    if kind == "bar":
        cats, vals = list(cats)[::-1], list(vals)[::-1]
    cd = CategoryChartData()
    cd.categories = list(cats)
    cd.add_series("금액", tuple(vals), number_format=num_fmt)

    ctype = XL_CHART_TYPE.COLUMN_CLUSTERED if kind == "col" else XL_CHART_TYPE.BAR_CLUSTERED
    gf = slide.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    ch.font.name = FONT
    ch.font.size = Pt(cat_size)
    ch.font.color.rgb = INK2

    plot = ch.plots[0]
    plot.gap_width = gap
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_fmt
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.name = FONT_M
    dl.font.size = Pt(label_size)
    dl.font.bold = True
    dl.font.color.rgb = INK

    ser = plot.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = color
    ser.format.line.fill.background()

    va = ch.value_axis
    va.visible = False
    va.has_major_gridlines = False
    va.has_minor_gridlines = False
    va.maximum_scale = max(vals) * headroom
    va.minimum_scale = 0

    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.major_tick_mark = XL_TICK_MARK.NONE
    ca.minor_tick_mark = XL_TICK_MARK.NONE
    ca.format.line.color.rgb = LINE
    ca.tick_labels.font.name = FONT
    ca.tick_labels.font.size = Pt(cat_size)
    ca.tick_labels.font.color.rgb = INK2

    cs = ch._chartSpace
    C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
    sp = parse_xml(
        '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>')
    cs.insert(list(cs).index(cs.find(C + "chart")) + 1, sp)
    return gf


# ── Ian 의 2026 CouSolve Idea Contest 제출본에서 가져온 표기 ──────────
# 그 자료가 이 양식을 실제로 어떻게 채웠는지가 기준이다:
#   · AS-IS / Pain Point, Target / Vision, Data / Issue  — 굵은 머리말 + 내용
#   · 'A vs. B'  ➡  우선한 쪽                            — 긴장 관계 표기
#   · 단계 흐름 (윗줄 굵게 / 아랫줄 기술 요소 이탤릭)
#   · 색 머리띠를 두른 KPI 카드
#   · 장 맨 아래 이탤릭 한 줄

def kv(slide, x, y, w, items, size=9.5, gap=0.30, key_color=INK):
    """'AS-IS : …' 처럼 굵은 머리말을 붙인 줄을 쌓는다."""
    for i, (k, v) in enumerate(items):
        text(slide, x, y + i * gap, w, gap,
             [("· ", {"color": INK3}), (k + " : ", {"bold": True, "color": key_color}),
              (v, {})], size=size)
    return y + len(items) * gap


def arrow(slide, x, y, w=0.42, h=0.20, color=RGBColor(0xB8, 0xBF, 0xC4)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def versus(slide, x, y, w, rows, size=9.5, gap=0.35, split=0.46):
    """'A vs. B'  ➡  우선한 쪽 — 계약사 자료의 Tradeoffs 표기."""
    lw = w * split
    for i, (left, right) in enumerate(rows):
        yy = y + i * gap
        text(slide, x, yy, lw - 0.30, gap,
             left if not isinstance(left, str) else [(left, {"bold": True})],
             size=size, align=PP_ALIGN.RIGHT)
        arrow(slide, x + lw - 0.22, yy + 0.05, 0.38, 0.17)
        text(slide, x + lw + 0.28, yy, w - lw - 0.28, gap,
             right if not isinstance(right, str) else [(right, {})], size=size)
    return y + len(rows) * gap


def kpi_card(slide, x, y, w, h, tag, name, body, color):
    """색 머리띠 + 흰 본문. 계약사 자료의 핵심 KPI 카드."""
    from pptx.enum.shapes import MSO_SHAPE
    head_h = 0.42
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(head_h))
    top.fill.solid(); top.fill.fore_color.rgb = color
    top.line.fill.background(); top.shadow.inherit = False
    top.text_frame.text = ""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + head_h),
                                 Inches(w), Inches(h - head_h))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = LINE; box.line.width = Pt(0.75); box.shadow.inherit = False
    box.text_frame.text = ""
    text(slide, x, y + 0.04, w, 0.20, [(tag, {"color": WHITE})], size=8, align=PP_ALIGN.CENTER)
    text(slide, x, y + 0.19, w, 0.24, [(name, {"bold": True, "color": WHITE})],
         size=10.5, align=PP_ALIGN.CENTER)
    yy = y + head_h + 0.12
    for chunks, sz in body:
        text(slide, x + 0.10, yy, w - 0.20, 0.24, chunks, size=sz, align=PP_ALIGN.CENTER)
        yy += 0.26
    return box


def quote(slide, y, t, size=9.5, x=0.85, w=11.70):
    return text(slide, x, y, w, 0.28, [("“" + t + "”", {"italic": True, "color": INK2})],
                size=size, align=PP_ALIGN.CENTER)


# 2계열 비교용 — dataviz 검증 통과 조합 (validate_palette.js, light, 6/6 PASS)
#   #1B8A6B ↔ #B9770B : CVD ΔE 8.9(protan) · 정상시 19.1 · 대비 3:1 이상
TEAL = RGBColor(0x1B, 0x8A, 0x6B)
AMBER = RGBColor(0xB9, 0x77, 0x0B)


def chart2(slide, x, y, w, h, cats, series, num_fmt="#,##0",
           label_size=9, cat_size=9, colors=(AMBER, TEAL), gap=80, headroom=1.20):
    """계열 두 개짜리 세로 막대. 계열이 둘이므로 범례를 반드시 둔다.

    series: [(이름, [값…]), (이름, [값…])]
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK,
                                 XL_LEGEND_POSITION)

    cd = CategoryChartData()
    cd.categories = list(cats)
    for name, vals in series:
        cd.add_series(name, tuple(vals), number_format=num_fmt)

    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.TOP
    ch.legend.include_in_layout = False
    ch.legend.font.name = FONT
    ch.legend.font.size = Pt(cat_size)
    ch.legend.font.color.rgb = INK2
    ch.font.name = FONT
    ch.font.size = Pt(cat_size)
    ch.font.color.rgb = INK2

    plot = ch.plots[0]
    plot.gap_width = gap
    plot.overlap = -12
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_fmt
    dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.name = FONT_M
    dl.font.size = Pt(label_size)
    dl.font.bold = True
    dl.font.color.rgb = INK

    for ser, col in zip(plot.series, colors):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = col
        ser.format.line.fill.background()

    va = ch.value_axis
    va.visible = False
    va.has_major_gridlines = False
    va.has_minor_gridlines = False
    va.maximum_scale = max(max(v) for _, v in series) * headroom
    va.minimum_scale = 0

    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.major_tick_mark = XL_TICK_MARK.NONE
    ca.minor_tick_mark = XL_TICK_MARK.NONE
    ca.format.line.color.rgb = LINE
    ca.tick_labels.font.name = FONT
    ca.tick_labels.font.size = Pt(cat_size)
    ca.tick_labels.font.color.rgb = INK2

    cs = ch._chartSpace
    C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
    sp = parse_xml(
        '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>')
    cs.insert(list(cs).index(cs.find(C + "chart")) + 1, sp)
    return gf
