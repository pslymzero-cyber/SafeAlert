# CouSolve Report 양식(7단계)에 SafeAlert 내용을 채운다.
#   python3 build_cousolve_deck.py <CouSolve_Report_양식.pptx>
#
# 양식 원본은 사내 문서라 저장소에 넣지 않는다 — 로컬 경로를 인자로 넘긴다.
# 양식의 12장 구성이 그대로 7단계라 장을 늘리거나 줄이지 않았다.
#   1 표지 / 2 1-1 Problem Situation / 3 1-2 Target / 4 1-3 Fishbone
#   5 1-4 Cause Analysis / 6 1-5 5Why / 7 2 Tradeoffs / 8 3 Benchmark
#   9 4 Solution / 10 5 Metrics / 11 6 Andon / 12 7 Feedback Loop
#
# 숫자 원칙 — 없는 데이터는 만들지 않는다.
# 사고 건수 · 셧다운 건수 · 절감 M/H 는 확보돼 있지 않다. 그 자리에는 '미확보' 라고 적는다.
# 금액은 절감액이 아니라 같은 범위를 하드웨어로 덮을 때의 환산액이다.
import os, shutil, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import _form as F

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(os.path.dirname(HERE), "company_form", "assets")
OUT = os.path.join(HERE, "SafeAlert_CouSolve_2026.08.pptx")

if len(sys.argv) < 2:
    raise SystemExit("사용법: python3 build_cousolve_deck.py <CouSolve_Report_양식.pptx>")
SRC = os.path.abspath(sys.argv[1])
BASE = os.path.join(HERE, "base_cousolve.pptx")
shutil.copyfile(SRC, BASE)
prs = Presentation(BASE)
assert len(prs.slides) == 12, f"양식이 12장이어야 한다 (지금 {len(prs.slides)}장)"


THEME = "SafeAlert  —  지급된 단말끼리 BLE 로 서로를 감지해, I-PAS 가 덮지 못하는 접근 조합을 경보한다"


def theme_bar(slide):
    """양식 상단의 빈 Theme 바(1.71, 0.87)에 주제를 한 줄 넣는다."""
    for sh in slide.shapes:
        if sh.left is None or not sh.has_text_frame:
            continue
        if abs(sh.left / 914400 - 1.71) < 0.02 and abs(sh.top / 914400 - 0.87) < 0.02:
            tf = sh.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            F._run(p, "  " + THEME, 10.5, False, F.INK2)
            return sh
    return None


def named(slide, *names):
    for sh in slide.shapes:
        if sh.name in names:
            return sh
    raise SystemExit(f"도형을 찾지 못했다: {names}")


def tables(slide):
    return [sh.table for sh in slide.shapes if sh.has_table]


def pic(slide, name, x, y, w, ratio, cap=None, cap_size=8.5):
    h = w / ratio
    slide.shapes.add_picture(os.path.join(ASSETS, name), Inches(x), Inches(y),
                             Inches(w), Inches(h))
    if cap:
        F.text(slide, x, y + h + 0.06, w, 0.24, [(cap, {"color": F.INK3})],
               size=cap_size, align=PP_ALIGN.CENTER)
    return y + h


# ── 1. 표지 ─────────────────────────────────────────────────
s = prs.slides[0]
sh = named(s, "직사각형 1")
tf = sh.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
for r in list(p.runs):
    r._r.getparent().remove(r._r)
F._run(p, "SafeAlert — 사각지대를 사람의 주의력에 맡기지 않는다", 20, True, F.GREEN)
p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
F._run(p2, "지급된 단말끼리 BLE 로 서로를 감지해, I-PAS 가 덮지 못하는 접근 조합을 메운다",
       11, False, F.INK2)
# 이 자료가 제안서와 갈리는 지점은 시제(時制)다. 표지에서 먼저 못박는다.
F.text(s, 2.41, 2.24, 8.54, 0.30,
       [("제안이 아니다  —  이미 현장에서 돌고 있다", {"bold": True, "color": F.GREEN})],
       size=13, align=PP_ALIGN.CENTER)
FACTS = [
    ("3개월 · 70회+", "현장에 배포하며 다듬은 기간"),
    ("3개 센터", "WF11 · WF21 · WF25 실사용 중"),
    ("리뷰 3차", "현장 요구 4건도 코드에 반영"),
    ("Phase 1 / 5", "신뢰성 로드맵 완료 단계"),
]
for i, (big, cap) in enumerate(FACTS):
    x = 2.41 + i * 2.135
    F.text(s, x, 2.82, 2.135, 0.36, [(big, {"bold": True, "color": F.RED})],
           size=15, align=PP_ALIGN.CENTER)
    F.text(s, x, 3.24, 2.135, 0.28, [(cap, {"color": F.INK2})], size=9, align=PP_ALIGN.CENTER)
F.text(s, 2.41, 3.68, 8.54, 0.28,
       [("추가 하드웨어 구매도, 외부 협조도 없이 지급된 단말만으로 동작한다",
         {"italic": True, "color": F.INK3})], size=9.5, align=PP_ALIGN.CENTER)

t = tables(s)[0]
F.cell(t.cell(0, 1), [("WF11  Waterflex", {})], 10, PP_ALIGN.CENTER)
F.cell(t.cell(0, 3), [("임효성 (Ian)", {})], 10, PP_ALIGN.CENTER)

for _s in list(prs.slides)[1:]:
    theme_bar(_s)

# ── 2. 1-1 Problem Situation ────────────────────────────────
s = prs.slides[1]
t = tables(s)[0]
def kvrow(items):
    out = []
    for k, v in items:
        out.append([("· ", {"color": F.INK3}), (k + " : ", {"bold": True}), (v, {})])
    return out

ROWS = [
    kvrow([("AS-IS", "존과 로케이션 사이 교차 구간에서 PIT 와 보행자가 서로를 보지 못한 채 접근"),
           ("Pain Point", "적재 파렛트가 시야를 가리고, 장비 소음에 경적이 묻혀 인지가 늦다")]),
    kvrow([("Target", "접근을 사전에 감지해 사람에게 알린다 — 태그 유무와 무관하게"),
           ("Vision", "I-PAS 가 3m 에서 장비를 세우기 전에, 15m · 8m 에서 사람이 먼저 안다")]),
    kvrow([("Data", "I-PAS 태그를 단 장비 · 인원만 덮인다. 보행자끼리 · EPJ · 미장착 장비는 대상 밖"),
           ("Issue", "그 구간에는 감지도 알림도 없어 위험 상황을 확인할 방법 자체가 없다")]),
    [[("덮이는 것은 태그를 산 만큼이다. 태그 없는 조합은 주의력에만 맡겨져 있고, "
       "대상이 늘 때마다 태그를 그만큼 더 사야 한다.", {})]],
    [[("Coverage : ", {"bold": True}),
      ("태그가 덮는 접근 조합은 ", {}), ("12쌍", {"bold": True, "color": F.RED}),
      (" (차량 4대 × 상시 3명). 인원이 늘어도 태그를 산 만큼만 는다.", {})],
     [("Cost : ", {"bold": True}),
      ("센터 1곳을 태그로 덮으면 3,310,000원. 17개 센터는 5,627만원 규모지만, "
       "확대 계획이 확정된 바 없어 절감액 · 회피액으로 주장하지 않는다.", {})]],
    [[("WF11 현장 확인 (VC OB 제외) · 단가 사용자 제공 2026.08 · 판정 반경 DevSettings.kt 실측값.",
       {})],
     [("경보 이력 : ", {"bold": True, "color": F.RED}),
      ("Firebase `alerts/{yyyyMMdd}` 에 등급 · 거리 · 시각이 일자별로 쌓이고 있다 "
       "(기기당 1분 1회). 집계하면 접근 조우 건수가 나온다 — 본 자료 시점 미집계.",
       {"color": F.RED})]],
]
for i, v in enumerate(ROWS):
    F.cell(t.cell(i, 1), v, 9.5, PP_ALIGN.LEFT, space=1.3)

pic(s, "fig_aisle.jpg", 7.83, 2.42, 4.70, 1200 / 750,
    "존과 로케이션 사이 통로 — 교차 구간에서 서로 보이지 않는다")
F.table(s, 7.83, 5.66, 4.70,
        ["접근 조합", "I-PAS", "SafeAlert"],
        [["PIT ↔ 보행자 (태그 착용)", "○", "○"],
         ["보행자 ↔ 보행자", "✕", "○"],
         ["EPJ ↔ 보행자", "✕", "○"],
         ["미장착 장비 ↔ 보행자", "✕", "○"]],
        col_w=[2.70, 1.00, 1.00], row_h=0.26, head_h=0.27, size=8.5, head_size=8.5,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER],
        colors=[F.INK, F.RED, F.GREEN])

# ── 3. 1-2 Target ───────────────────────────────────────────
s = prs.slides[2]
t = tables(s)[0]
TARGET = [
    [[("I-PAS 가 덮지 못하는 접근 조합에서, PIT · EPJ · 보행자가 서로에게 다가오는 것을 "
       "사전에 감지해 경보한다.", {"bold": True})],
     [("장비를 세우는 일은 하지 않는다 — 3m 자동 셧다운은 I-PAS EOD 가 그대로 맡는다.",
       {"color": F.INK2})]],
    [[("이미 지급된 단말에 앱을 설치해 BLE 로 서로를 감지한다. 앵커 · 배선 · 서버가 없다.", {})],
     [("신호 수신  →  거리 추정  →  등급 판정  →  3중 경보", {"bold": True})],
     [("BLE 광고   →   3단 필터   →   상태머신   →   소리 · 진동 · 화면",
       {"italic": True, "color": F.INK3, "size": 9})],
     [("판정 반경 — 지게차가 낀 조합 경고 15m / 위험 8m,  그 밖의 조합 경고 5m / 위험 3m.",
       {"color": F.INK2})]],
    [[("추가 하드웨어 ", {}), ("0원", {"bold": True, "color": F.RED}),
      ("으로 센터 1곳 3,310,000원 범위를 덮는다. 17개 센터 확산 시 56,270,000원 규모.", {})],
     [("절감액이 아니라 같은 범위를 하드웨어로 살 때의 환산액이다. 앱은 그 구매를 만들지 않는다.",
       {"color": F.INK2})]],
    [[("Phase 1 (테스트 하네스 · CI 회귀 게이트) 완료. Phase 2 (골든 테스트 + PDA 이식) 착수 "
       "결정이 필요한 시점이다.", {})],
     [("Phase 3~5 는 순차. 결정이 늦어도 v1.1.70 이 현장에서 계속 돈다 — 보호가 멈추지는 않는다.",
       {"color": F.INK2})]],
]
for i, v in enumerate(TARGET):
    F.cell(t.cell(i, 1), v, 9.5, PP_ALIGN.LEFT, space=1.3)

# ── 4. 1-3 Fishbone Diagram ─────────────────────────────────
# 양식의 뼈대는 그룹 도형이라 건드리지 않고, 가시선 위에 글자만 얹는다.
# 그룹 좌표계(chOff/chExt) → 슬라이드 좌표 변환은 아래 두 줄이 전부다.
s = prs.slides[3]
GX, GY, SX, SY = 0.8382, 2.2092, 0.94937, 0.97129
CHX, CHY = 0.5633, 1.7510
to_x = lambda cx: GX + (cx - CHX) * SX
to_y = lambda cy: GY + (cy - CHY) * SY
BONE_W = 2.27 * SX

BONES = [
    # (뼈 왼쪽 끝 child x, child y, 항목 3개)
    (1.82, [2.70, 3.16, 3.62], ["육안 · 경적에만 의존", "위험 상황 기록 체계 없음",
                                "안전 설비가 장비 단위로만 설계"]),          # Method
    (5.21, [2.70, 3.16, 3.62], ["적재 파렛트에 시야 가림", "장비 소음에 경적이 묻힘",
                                "단독 작업 구간 상호 확인 불가"]),           # Man
    (1.82, [4.81, 5.27, 5.73], ["I-PAS 는 태그 대상만 감지", "EOD 는 3m 에서만 개입",
                                "EPJ · 미장착 장비는 대상 밖"]),             # Machine
    (5.21, [4.81, 5.27, 5.73], ["태그 1세트 855,000원", "보행자 태그 = 상시 인원분",
                                "비콘 · 거치대 설치가 별도"]),               # Material
]
for cx, cys, items in BONES:
    for cy, it in zip(cys, items):
        F.text(s, to_x(cx) + 0.04, to_y(cy) - 0.235, BONE_W - 0.06, 0.22,
               [(it, {})], size=8.5, color=F.INK2)

F.text(s, to_x(10.14) + 0.10, to_y(3.30) + 0.28, 2.354 - 0.20, 1.30,
       [("사각지대에서 서로를 보지 못한 채 접근한다", {"bold": True, "color": F.GREEN})],
       size=10.5, align=PP_ALIGN.CENTER)
F.text(s, to_x(10.14) + 0.10, to_y(3.30) + 1.00, 2.354 - 0.20, 0.70,
       [("Machine · Material 이 같은 뿌리를 가리킨다 — "
         "감지 주체가 별도 하드웨어에 묶여 있다", {"color": F.GREEN2})],
       size=8.5, align=PP_ALIGN.CENTER, space=1.25)

# ── 5. 1-4 Cause Analysis ───────────────────────────────────
s = prs.slides[4]
t = tables(s)[0]
CAUSES = [
    ("적재 파렛트가 시야를 가린다", "WF11 통로 현장 확인", "확인됨 — 교차 구간 상호 시야 불가"),
    ("장비 소음으로 경적 인지가 늦다", "현장 청취 · 작업자 진술", "확인됨"),
    ("I-PAS 는 태그 장착 대상만 덮는다", "I-PAS 구성 확인 (차량 태그 · EOD · 보행자 태그)",
     "확인됨 — 보행자끼리 · EPJ 는 대상 밖"),
    ("대상이 늘면 태그를 그만큼 더 산다", "단가 확인 (1세트 855,000원, 2026.08)",
     "확인됨 — 확산이 곧 구매"),
    ("위험 상황이 집계되지 않는다", "안전 시스템 유무 확인",
     [[("확인됨 — ", {}), ("사고 · 셧다운 건수 데이터 없음", {"color": F.RED})]]),
    ("경보가 잦으면 작업자가 꺼버린다", "리더십 리뷰 2차 피드백",
     "확인됨 → 자동 음소거 · 세이프존으로 대응"),
]
for i, row in enumerate(CAUSES, start=1):
    for j, v in enumerate(row):
        F.cell(t.cell(i, j), v if not isinstance(v, str) else [(v, {})],
               9.5, PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, space=1.15)
F.cell(t.cell(7, 0),
       [[("확인 방법은 현장 확인 · 단가 대조 · 배포 이력 세 가지뿐이다. ", {}),
         ("계측 데이터로 검증한 항목은 없다 — 감지 수단이 없어 측정된 적이 없기 때문이다.",
          {"color": F.RED})]],
       9, PP_ALIGN.LEFT)

# ── 6. 1-5 5Why ─────────────────────────────────────────────
s = prs.slides[5]
t = tables(s)[0]
F.cell(t.cell(0, 0),
       [[("* Most likely causes :   ", {"color": F.INK2}),
         ("감지 대상이 태그를 단 장비와 태그를 받은 인원으로만 정의돼 있다", {"bold": True})]],
       10, PP_ALIGN.LEFT)
WHY = [
    (1, "왜 사각지대에서 위험한가?", "서로를 보지 못한 채 접근한다"),
    (3, "왜 보지 못하나?", "적재 파렛트가 시야를 가리고, 장비 소음에 경적이 묻힌다"),
    (5, "왜 그 구간에 보완 수단이 없나?",
     "안전 설비를 장비 단위로 도입해 왔다 — 차량에 리더를 달고, 그 차량이 감지한다"),
    (7, "왜 장비 단위인가?",
     "감지 대상이 ‘차량’ 으로 정의돼 있고, 사람은 태그를 받은 만큼만 대상이 된다"),
    (9, "왜 인원 전원이 대상이 아닌가?",
     "인원 수만큼 태그를 사야 하는 구조라, 대상을 상시 인원 3명으로 한정했다"),
]
for r, q, ans in WHY:
    F.cell(t.cell(r, 1), [(q, {"color": F.GREEN2, "bold": True})], 10, PP_ALIGN.LEFT)
    F.cell(t.cell(r + 1, 1), [[("→   ", {"color": F.INK3}), (ans, {})]], 9.5, PP_ALIGN.LEFT)
F.cell(t.cell(11, 0),
       [[("* Root Cause :   ", {"color": F.INK2}),
         ("감지 대상 범위가 안전 요구가 아니라 태그 구매 수량으로 결정되고 있다",
          {"bold": True, "color": F.RED}),
         ("   (외부 상용 제품도 같은 구조 — 3장 참조)", {"color": F.INK3})]],
       10, PP_ALIGN.LEFT)

# ── 7. 2. Tradeoffs ─────────────────────────────────────────
# 계약사 자료의 'A vs. B → 우선한 쪽' 표기를 그대로 쓴다.
# 위 구역은 무엇을 포기했는지, 아래 구역은 포기한 쪽을 무엇으로 메웠는지다.
s = prs.slides[6]
F.versus(s, 2.25, 2.28, 10.35, [
    ("감지 민감도  vs.  경보 피로", [("민감도를 택했다", {"bold": True}),
      ("  —  넓게 울리되, 울릴 곳을 좁힌다", {"color": F.INK2})]),
    ("감지 주기  vs.  배터리 · 발열", [("1교대 구동을 상한으로 둔다", {"bold": True}),
      ("  —  주기를 그 안에서만 촘촘히", {"color": F.INK2})]),
    ("배포 속도  vs.  판정 재현성", [("재현성을 택했다", {"bold": True}),
      ("  —  3개월 70회에서 게이트 통과 방식으로", {"color": F.INK2})]),
    ("앱 경보  vs.  물리 정지", [("물리 정지는 포기했다", {"bold": True, "color": F.RED}),
      ("  —  I-PAS EOD 를 그대로 남긴다", {"color": F.INK2})]),
    ("개인 폰 임시 설치  vs.  지속 가능성", [("지속 가능성을 택했다", {"bold": True}),
      ("  —  업무 단말로 옮긴다", {"color": F.INK2})]),
], size=9.5, gap=0.345)

F.zone(s, 2.25, 4.20, 10.35, [
    (None, [
        [("자동 음소거", {"bold": True, "color": F.GREEN2}),
         ("   같은 기기 · 같은 등급에 5초 이상 머물면 소리 · 진동을 끈다. 표시 · 기록은 유지하고, "
          "등급이 오르면 즉시 재발령", {}), ("  (v1.1.61)", {"color": F.INK3})],
        [("세이프존", {"bold": True, "color": F.GREEN2}),
         ("   휴게실 · 충전 구역에 비콘을 두면 그 안에서는 경보를 억제한다. 존 안의 기기는 "
          "상대에게도 안전으로 보인다", {}), ("  (v1.1.62)", {"color": F.INK3})],
        [("이탈 재발령 억제", {"bold": True, "color": F.GREEN2}),
         ("   멀어진 뒤 경고 범위에서 다시 울리지 않게 했다", {}), ("  (v1.1.51)", {"color": F.INK3})],
        [("CI 회귀 게이트", {"bold": True, "color": F.GREEN2}),
         ("   고친 증상을 테스트로 고정한다. 골든 테스트가 실패하면 릴리스가 자동 차단된다", {}),
         ("  (Phase 1 완료)", {"color": F.INK3})],
        [("통제 위계", {"bold": True, "color": F.GREEN2}),
         ("   자동 정지(공학적 제어)가 상위, 경보(관리적 제어)가 하위다. 상위를 대체하지 않는다 — "
          "I-PAS EOD 를 그대로 두고, 상위가 닿지 않는 조합에만 하위를 새로 넣는다", {})],
        [("PDA 이식", {"bold": True, "color": F.GREEN2}),
         ("   Phase 2 에서 업무 단말로 옮긴다. 개인 폰 의존과 배터리 부담이 함께 없어진다", {}),
         ("  (예정)", {"color": F.INK3})],
    ]),
], gap=0.375)
F.quote(s, 6.66, "완벽한 판정보다 현장이 계속 켜 두는 경보를 택했다. 꺼 둔 단말은 보호가 0 이다.",
        x=2.25, w=10.35)

# ── 8. 3. Benchmark ─────────────────────────────────────────
# 양식이 요구하는 벤치마크는 '다른 회사는 어떻게 해결했는가' 다.
# 사내 장비 비교만 넣으면 양식 미준수다. 외부 상용 시스템을 먼저 놓고 그 안에 I-PAS 를 둔다.
s = prs.slides[7]
F.text(s, 2.25, 2.26, 10.35, 0.26,
       [("지게차 ↔ 보행자 근접 경보는 이미 상용 제품군이 있다. 우리 현행 장비(I-PAS)도 그중 하나다.",
         {"color": F.INK2})], size=9.5)
F.table(s, 2.25, 2.56, 10.35,
        ["시스템", "방식", "감지 거리", "착용물", "차량 개입"],
        [["ZoneSafe", "RFID", "차량 주위 360° 최대 10m", "태그 필요", "경보"],
         ["ELOKON ELOshield", "UWB", "경고 / 보호 구역 분리 · 비가시선 감지", "태그 필요", "경보 + 자동 감속"],
         ["Pozyx RTLS", "UWB", "차량 주위 360° 최대 25m", "태그 필요", "경보"],
         ["I-PAS  (현행)", "태그", "3m", "태그 필요", "경보 + 자동 셧다운"],
         [[("SafeAlert", {"bold": True, "color": F.TEAL})], "BLE", "경고 15m / 위험 8m",
          [[("없음 — 지급 단말", {"bold": True, "color": F.TEAL})]], "경보만"]],
        col_w=[2.30, 0.95, 3.65, 1.85, 1.60], row_h=0.28, head_h=0.28,
        size=8.5, head_size=8.5,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])

F.lines(s, 2.35, 4.22, 10.25, [
    [("·  외부 제품은 예외 없이 ", {}), ("‘차량에 리더, 사람에 태그’", {"bold": True}),
     (" 구조다. 대상이 늘면 태그를 산다 — I-PAS 와 같다.", {})],
    [("·  감지 거리 10~25m 는 SafeAlert 의 경고 15m 와 같은 범위다. 반경 설정이 업계 관행에서 "
      "벗어나 있지 않다.", {})],
    [("·  상위 제품은 자동 감속까지 간다. 그 자리는 I-PAS EOD 가 이미 맡고 있고, "
      "SafeAlert 는 그 앞단만 맡는다.", {"color": F.RED})],
    [("·  다른 점은 ", {}), ("착용물의 유무 하나", {"bold": True}),
     ("다. 대상 확대 비용의 차이가 전부 거기서 나온다 — 아래가 그 차이다.", {})],
], size=9, gap=0.255)
F.text(s, 2.35, 5.26, 10.25, 0.22,
       [("벤더 공개 자료 기준 (2026.09 조회). ‘아차사고 감소 · ROI 3~6개월’ 은 벤더 주장이며 "
         "검증된 수치가 아니라 인용하지 않는다.", {"color": F.INK3})], size=8)

# 인원이 늘 때 덮이는 조합 수 — 태그는 선형(V×T), 앱은 제곱(C(V+P,2)).
# 이건 산술이라 반박 대상이 아니다. '안 쓸 돈' 을 절감액처럼 말하지 않아도 된다.
V = 4
PEOPLE = [3, 5, 10, 20]
tag_pairs = [V * p for p in PEOPLE]
app_pairs = [(V + p) * (V + p - 1) // 2 for p in PEOPLE]
F.chart2(s, 2.35, 5.44, 10.25, 1.62,
         [f"현장 인원 {p}명" for p in PEOPLE],
         [("태그 방식 — 태그 착용자만 (차량 4대 기준, 인원당 110,000원)", tag_pairs),
          ("SafeAlert — 앱 설치 단말 전부 (추가 비용 0원)", app_pairs)],
         label_size=8.5, cat_size=8.5)

# ── 9. 4. Solution ──────────────────────────────────────────
s = prs.slides[8]
left = [t for t in tables(s) if len(t.columns) == 1][0]
plan = [t for t in tables(s) if len(t.columns) == 4][0]
F.cell(left.cell(1, 0),
       [[("지급 단말에 SafeAlert 를 설치해 BLE 로 상호 감지한다.", {"bold": True})],
        [("지게차가 낀 조합 경고 15m / 위험 8m, 그 밖의 조합 5m / 3m. 소리 · 진동 · 화면 3중 경보. "
          "v1.1.70 이 WF11 · WF21 · WF25 에서 운영 중.", {"color": F.INK2})]],
       9.5, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, space=1.25)
F.cell(left.cell(3, 0),
       [("추가 하드웨어 없이 I-PAS 가 비워 둔 조합을 지금 덮는다. 대상이 늘어도 구매가 없다.", {})],
       9.5, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, space=1.2)
F.cell(left.cell(5, 0),
       [[("신뢰성 로드맵 Phase 1~5 + Phase 2 PDA 이식.", {"bold": True})],
        [("새 기능을 붙이는 작업이 아니다. 판정이 흔들리는 구조적 원인을 걷어내고, "
          "개인 폰 임시 설치를 업무 단말로 옮긴다.", {"color": F.INK2})]],
       9.5, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, space=1.25)
F.cell(left.cell(7, 0),
       [("판정 재현성이 확보돼야 확산할 수 있다. PDA 로 옮기면 개인 폰 의존과 배터리 부담이 함께 없어진다.",
         {})], 9.5, PP_ALIGN.LEFT, MSO_ANCHOR.TOP, space=1.2)

PLAN = [
    ("Phase 1  테스트 하네스 · CI 회귀 게이트", "개발", "완료", "완료"),
    ("경보 이력 집계 — Firebase alerts, 4주분", "개발", "즉시 가능", "착수"),
    ("3개 센터 실사용 테스트 승인 근거 확인", "안전 · 리더십", "제출 전", "확인 필요"),
    ("보안 검토 — 난독화 미적용 · 이력 평문 저장", "보안", "확산 전", "요청"),
    ("서베이 회수 · 집계 (WF11 · WF21 · WF25)", "현장", "진행 중", "진행"),
    ("Phase 2  안전 경로 골든 테스트 + PDA 이식", "개발 · 현장", "리소스 배정 후", "예정"),
    ("Phase 3~5  분해 · 상태 단일화 · 워커 분리", "개발", "Phase 2 후", "예정"),
    ("정식 과제 등록 · 유지보수 주체 지정", "리더십", "결정 필요", "요청"),
    ("현장 검증 슬롯 (단계당 확인 1건)", "현장", "단계별", "요청"),
    ("17개 센터 확산", "리더십 · 현장", "Phase 2 완료 후", "미정"),
]
STCOL = {"완료": F.GREEN, "착수": F.TEAL, "확인 필요": F.RED, "예정": F.INK2,
         "요청": F.RED, "진행": F.GREEN2, "미정": F.INK3}
for i, (what, who, when, st) in enumerate(PLAN, start=2):
    F.cell(plan.cell(i, 0), [(what, {})], 9, PP_ALIGN.LEFT, margin=0.06)
    F.cell(plan.cell(i, 1), [(who, {})], 9, PP_ALIGN.CENTER, margin=0.03)
    F.cell(plan.cell(i, 2), [(when, {})], 9, PP_ALIGN.CENTER, margin=0.03)
    F.cell(plan.cell(i, 3), [(st, {"bold": True, "color": STCOL[st]})], 9, PP_ALIGN.CENTER, margin=0.03)

# ── 10. 5. Metrics ──────────────────────────────────────────
# 「미발령 0건」은 뺐다 — 미발령은 정의상 관측되지 않아 지표로 성립하지 않는다.
# 카드 색이 곧 준비 상태다: 청록 = 지금 되는 것, 호박 = 확보 필요.
s = prs.slides[9]
KPI = [
    ("KPI 1", "접근 조우", F.TEAL, "위험 · 경고 등급 발령 건수", "Firebase 실측", "지금 집계 가능"),
    ("KPI 2", "회귀 재발", F.TEAL, "고친 증상이 다시 관측", "CI 골든 테스트", "측정 중 · 목표 0건"),
    ("KPI 3", "오발령 신고", F.AMBER, "정지 중 반복 · 이탈 후 재발령", "작업자 신고", "채널 신설 필요"),
    ("KPI 4", "체감 만족도", F.AMBER, "경보 타이밍 · 거리 정확도", "10점 척도 서베이", "회수 필요"),
]
for i, (tag, name, col, desc, how, state) in enumerate(KPI):
    x = 2.35 + i * 2.56
    F.kpi_card(s, x, 2.26, 2.42, 1.62, tag, name,
               [([(desc, {"color": F.INK2})], 8.5),
                ([(how, {"color": F.INK3, "italic": True})], 8),
                ([(state, {"bold": True, "color": col})], 9.5)],
               col)
F.lines(s, 2.35, 3.96, 10.25, [
    [("「미발령 0건」은 지표에서 뺐다", {"bold": True, "color": F.RED}),
     ("  —  미발령은 정의상 관측되지 않는다. 서베이의 ‘경보가 늦었다’ 응답으로 간접 관측한다.", {})],
], size=9.5, gap=0.28)

F.lines(s, 2.35, 4.30, 10.25, [
    [("·  KPI 1 은 지금 뽑을 수 있다 — ", {}),
     ("`alerts/{yyyyMMdd}/{id}` = timestamp · deviceId · walkerId · rssi · alertLevel",
      {"bold": True}),
     (" 이 이미 쌓여 있다 (FirebaseManager.kt:15-29).", {})],
    [("·  기기당 1분 1회로 스로틀돼 있어(BleService.kt:2529) 중복이 제거된 상태다 — "
      "건수가 곧 ‘접근 조우 횟수’ 에 가깝다. 아차사고 대리지표로 쓸 수 있다.", {})],
    [("·  KPI 2 는 안전 크리티컬 경로 골든 테스트가 실패하면 릴리스가 자동 차단되는 형태로 "
      "이미 동작 중이다 (Phase 1).", {})],
    [("·  KPI 3 · 4 는 아직 없다. 신고 채널과 서베이 회수가 선행되어야 한다.", {"color": F.AMBER})],
    [("·  사고 건수 · 셧다운 건수는 회사 데이터로 존재하지 않는다. ", {}),
     ("추정치를 만들지 않고, 지표에서 제외했다.", {"color": F.RED})],
], size=9.5, gap=0.315)
F.quote(s, 6.30, "세는 방법이 없는 것은 지표로 쓰지 않는다. 대신 이미 세고 있는 것부터 꺼낸다.",
        x=2.25, w=10.35)

# ── 11. 6. Andon ────────────────────────────────────────────
# 개발 게이트가 아니라 현장에서 실제로 벌어지는 상황과, 업데이트 · 수정 중 상황이다.
s = prs.slides[10]
STOP, WATCH = "즉시 중단", "관측 시 판단"
SCEN = [
    ("현장 운영 중", [
        ("미발령", STOP, "접근했는데 경보가 뜨지 않은 사례 1건이라도"),
        ("오발령 지속", STOP, "정지 중 계속 울림 · 이탈 후 다시 울림이 반복"),
        ("감지 중단", WATCH, "상시 알림이 ‘이상’ 으로 바뀌거나 사라짐"),
        ("경보 피로", WATCH, "무음으로 두거나 앱을 꺼 둔 단말이 늘어남"),
    ]),
    ("업데이트 · 수정 중", [
        ("버전 혼재", STOP, "구버전 단말이 신버전을 감지하지 못함"),
        ("회귀 재발", STOP, "이전에 고친 증상이 다시 관측됨"),
        ("보호 공백", WATCH, "작업 시간 중 일괄 업데이트 — 교대 전환 때만 배포"),
        ("서비스 미기동", WATCH, "업데이트 · 역할 전환 후 백그라운드 실행 미복구"),
    ]),
]
for ci, (head, items) in enumerate(SCEN):
    cx = 2.25 + ci * 5.20
    F.text(s, cx, 2.26, 5.00, 0.26,
           [("▪ ", {"color": F.GREEN2, "bold": True}), (head, {"bold": True, "color": F.GREEN2})],
           size=10)
    for i, (name, tag, desc) in enumerate(items):
        y = 2.58 + i * 0.40
        F.text(s, cx + 0.10, y, 4.85, 0.22,
               [(name, {"bold": True}), ("   " + tag, {"bold": True, "size": 8.5,
                                                       "color": F.RED if tag == STOP else F.INK2})],
               size=9.5)
        F.text(s, cx + 0.10, y + 0.19, 4.85, 0.21, [(desc, {"color": F.INK2})], size=8.5)

ROLL = [
    ("즉시 롤백", "직전 태그 APK 로 되돌린다 — 현장 재설치 없이 배포된다"),
    ("선행 확인", "전 단말 배포 전 1대에서 먼저 확인한다"),
    ("빌드 차단", "골든 테스트가 실패하면 릴리스가 자동 차단된다 (동작 중)"),
    ("원인 귀속", "단계를 단독 배포하므로 어느 변경 탓인지 특정된다"),
    ("재개 조건", "실패한 상황을 테스트로 고정한 뒤에만 다음 단계로"),
]
for i, (k, v) in enumerate(ROLL):
    F.text(s, 2.25, 4.24 + i * 0.34, 10.35, 0.30,
           [(k, {"bold": True, "color": F.GREEN2}), ("   —   ", {"color": F.INK3}), (v, {})],
           size=9.5)
F.text(s, 2.25, 6.06, 10.35, 0.28,
       [("멈출 기준을 먼저 정해 두는 것이, 멈추지 않고 밀어붙이는 것보다 빠르다.",
         {"italic": True, "color": F.INK2})], size=9.5)

# ── 12. 7. Feedback Loop ────────────────────────────────────
s = prs.slides[11]
F.text(s, 0.85, 2.16, 5.80, 0.26,
       [("▪ ", {"color": F.GREEN2, "bold": True}), ("리더십 리뷰 반영", {"bold": True, "color": F.GREEN2})],
       size=10)
F.table(s, 0.85, 2.48, 5.75,
        ["차수", "반영 내용", "결과"],
        [["1차", "역할 전환 버튼 · EPJ 역할 숨김 · 스플래시 교체", "재설치 없이 화면에서 역할 전환"],
         ["2차", "같은 기기 · 같은 등급 5초 이상이면 자동 음소거", "표시 · 기록 유지, 등급 오르면 재발령"],
         ["3차", "세이프존 — 휴게실 · 충전 구역 비콘으로 경보 억제", "존 안의 기기는 상대에게도 안전"]],
        col_w=[0.60, 2.95, 2.20], row_h=0.44, head_h=0.28, size=8.5, head_size=8.5,
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
F.text(s, 6.90, 2.16, 5.80, 0.26,
       [("▪ ", {"color": F.GREEN2, "bold": True}), ("현장 요구 반영", {"bold": True, "color": F.GREEN2})],
       size=10)
F.table(s, 6.90, 2.48, 5.75,
        ["항목", "Before", "After"],
        [["보호 상태", "끊겨도 작업자가 알 수 없음", "감지가 멈추면 상시 알림에 이상 표시"],
         ["이탈 후", "멀어졌는데 다시 울림", "경고 범위 재발령 억제"],
         ["알림 끄기", "단계가 번거로움", "알림 본문을 누르면 즉시 무음"],
         ["경보 화면", "눈에 잘 들어오지 않음", "화면 가장자리에 붙는 사이드바"]],
        col_w=[1.05, 2.30, 2.40], row_h=0.32, head_h=0.28, size=8.5, head_size=8.5)
F.text(s, 0.85, 4.62, 11.70, 0.26,
       [("반영은 코드에 남아 있다  —  ", {"bold": True, "color": F.GREEN2}),
        ("리뷰 1·2·3차 v1.1.60 · v1.1.61 · v1.1.62,  보호 끊김 v1.1.64,  이탈 재발령 v1.1.51,  "
         "알림 본문 탭 v1.1.68,  사이드바 v1.1.69 · v1.1.70", {})], size=9)
F.text(s, 0.85, 4.96, 11.70, 0.60,
       [("피드백이 들어오면 그 주에 배포했다. 3개월 70회 이상. 문제는 속도가 아니라 재현성이었다  —  "
         "빨리 고친 증상이 되돌아왔고, 그래서 다음 작업은 새 기능이 아니라 테스트와 CI 게이트다.", {})],
       size=9.5, space=1.3)

F.lines(s, 0.85, 6.46, 5.70, [
    "역할 전환 · 자동 음소거 · 세이프존으로 경보 피로를 억제했다",
    "보호 끊김과 서비스 미기동이 상시 알림에 드러난다",
    "Phase 1 완료 — 고친 증상은 테스트로 고정돼 빌드에서 걸린다",
], size=8.5, gap=0.215)
F.lines(s, 6.94, 6.46, 5.70, [
    "WF11 · WF21 · WF25 서베이 회수 · 집계 (10점 척도)",
    "Phase 2 PDA 이식 — 개인 폰 임시 설치를 업무 단말로",
    "정식 과제 등록 · 리소스 배정이 선행되어야 17개 센터로 확산",
], size=8.5, gap=0.215)

prs.save(OUT)
print("wrote", OUT)
