# CouSolve Idea Contest 자료에서 가져온 표현 장치들.
# 사내 폼(LFP)의 흰 배경·제목 자리는 그대로 두고, 슬라이드 안쪽만 이 장치로 바꾼다.
#   pill      — 회색 라벨 (문제의 배경 / As-is / To-be / Tradeoffs …)
#   kv        — 굵은 라벨 + 설명 한 줄 (AS-IS : … / Pain Point : …)
#   table     — 어두운 머리행 + 가는 밑줄 표 (As-is·To-be, Before·After, 단계)
#   arrow     — 트레이드오프 좌우를 잇는 굵은 회색 화살표
#   headline  — 번호 매긴 한 줄 요약
#   quote     — 슬라이드 맨 아래 이탤릭 결론 한 문장
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

INK   = RGBColor(0x1C, 0x28, 0x33)
INK2  = RGBColor(0x5D, 0x6D, 0x7E)
INK3  = RGBColor(0x8A, 0x94, 0x9E)
GRAY  = RGBColor(0x8C, 0x8C, 0x8C)   # pill 배경
DARK  = RGBColor(0x33, 0x3B, 0x45)   # 표 머리행
LINE  = RGBColor(0xD5, 0xDB, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xC0, 0x39, 0x2B)
FONT  = "Malgun Gothic"


def _tf(shape, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(para, text, size, bold=False, color=INK, italic=False):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return r


def text(slide, x, y, w, h, chunks, size=12, align=PP_ALIGN.LEFT, space=4):
    """chunks: [(글자, {bold/color/italic/size})] 또는 [[...], [...]] 로 여러 문단."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = _tf(box)
    paras = chunks if chunks and isinstance(chunks[0], list) else [chunks]
    for i, para_chunks in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for t, opt in para_chunks:
            _run(p, t, opt.get("size", size), opt.get("bold", False),
                 opt.get("color", INK), opt.get("italic", False))
    return box


def pill(slide, x, y, label, w=1.55, h=0.33, size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = GRAY
    sh.line.fill.background()
    sh.shadow.inherit = False
    p = _tf(sh).paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, label, size, True, WHITE)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sh


def kv(slide, x, y, w, rows, size=12, gap=0.34):
    """[(라벨, 내용)] — 라벨만 굵게. CouSolve 의 'AS-IS : …' 두 줄 표기."""
    for i, (k, v) in enumerate(rows):
        text(slide, x, y + i * gap, w, gap, [(k + " : ", {"bold": True}), (v, {})], size=size)


def arrow(slide, x, y, w=0.9, h=0.42):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xC7, 0xCC, 0xD1)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _border(cell, edge, color=LINE, pt=0.75):
    tag = {"L": "a:lnL", "R": "a:lnR", "T": "a:lnT", "B": "a:lnB"}[edge]
    tcPr = cell._tc.get_or_add_tcPr()
    for e in tcPr.findall(qn(tag)):
        tcPr.remove(e)
    ln = parse_xml(
        '<{t} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'w="{w}" cap="flat" cmpd="sng" algn="ctr">'
        '<a:solidFill><a:srgbClr val="{c}"/></a:solidFill></{t}>'
        .format(t=tag, w=int(pt * 12700), c="%02X%02X%02X" % (color[0], color[1], color[2]))
    )
    tcPr.insert(0, ln)


def table(slide, x, y, w, header, rows, col_w, row_h=0.38, head_h=0.36, size=11, head_size=11,
          aligns=None, colors=None):
    """어두운 머리행 + 가는 밑줄. col_w 는 인치 비율이 아니라 실제 인치."""
    n_r, n_c = len(rows) + 1, len(header)
    gf = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(head_h + row_h * len(rows)))
    tbl = gf.table
    # 기본 표 스타일(줄무늬 파랑)을 끈다
    tbl._tbl.find(qn("a:tblPr")).set("bandRow", "0")
    tbl._tbl.find(qn("a:tblPr")).set("firstRow", "1")
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r):
        tbl.rows[i].height = Inches(row_h)

    for j, htxt in enumerate(header):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = DARK
        c.margin_left = c.margin_right = Inches(0.08)
        c.margin_top = c.margin_bottom = 0
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, htxt, head_size, True, WHITE)
        for e in "LRTB":
            _border(c, e, DARK)

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = (aligns or [PP_ALIGN.LEFT] * n_c)[j]
            col = (colors or [INK] * n_c)[j]
            bold = j == 0 and n_c > 2
            _run(p, val, size, bold, col)
            for e in "LRT":
                _border(c, e, WHITE, 0.75)
            _border(c, "B", LINE, 0.75)
    return gf


def headline(slide, y, parts, size=15, x=0.7, w=11.93):
    """'1. 사각지대 해소  2. 인프라 없는 확산 …' — 번호만 붉게."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.42))
    p = _tf(box).paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for i, t in enumerate(parts, start=1):
        _run(p, f"{i}. ", size, True, RED)
        _run(p, t + ("      " if i < len(parts) else ""), size, True, INK)
    return box


def quote(slide, y, t, size=11.5, x=0.7, w=11.93):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.4))
    p = _tf(box).paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "“" + t + "”", size, False, INK2, italic=True)
    return box


def flow(slide, x, y, w, steps, subs, size=14, sub_size=10.5):
    """'출발 정보 → 도착 예측 → …' 아래에 기술 요소를 작게 붙이는 CouSolve Solution 표기."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.7))
    tf = _tf(box)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for i, s in enumerate(steps):
        _run(p, s, size, True, INK)
        if i < len(steps) - 1:
            _run(p, "  →  ", size, True, INK3)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    for i, s in enumerate(subs):
        _run(p2, s, sub_size, False, INK3, italic=True)
        if i < len(subs) - 1:
            _run(p2, "   →   ", sub_size, False, INK3, italic=True)
    return box


def clear_body(shape):
    """본문 텍스트박스를 비운다 (도형으로 대체할 슬라이드용)."""
    tf = shape.text_frame
    txBody = tf.paragraphs[0]._p.getparent()
    NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for p in list(txBody.findall(NS + "p"))[1:]:
        txBody.remove(p)
    for r in list(tf.paragraphs[0].runs):
        r._r.getparent().remove(r._r)
