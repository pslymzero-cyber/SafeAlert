# LFP 사내 폼(표지·목차·본문 레이아웃)을 그대로 쓰고 내용만 SafeAlert 로 갈아 끼운다.
# 본문 장은 CouSolve Idea Contest 자료의 표현 장치(_devices.py)로 바꿔 그렸다 —
# 회색 라벨 · As-is/To-be 표 · 비교표 · 하단 이탤릭 결론.
# 구성은 매니저 지시(2026.08) 5장 — 개발배경 / 기대효과(사각지대 · 비용) / 장단점 /
# 앞으로 수정계획 / 테스트 중간 결과. 지시에 없던 장은 부록으로 내렸다.
# 폼 자체(흰 배경 · 상단 색 바 · 로고 · CFS 쪽번호)는 사내 표준 그대로 둔다.
import copy, os
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

import _devices as D

SRC, OUT = "base.pptx", "SafeAlert_보고_2026.08.pptx"
HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
A = lambda n: f"{HERE}/assets/{n}"

prs = Presentation(SRC)
NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def body_of(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and not sh.is_placeholder:
            return sh
    raise SystemExit("본문 텍스트박스를 찾지 못했다")


def title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.has_text_frame and "번호" not in sh.name:
            return sh
    raise SystemExit("제목 자리표시자를 찾지 못했다")


def set_text(shape, text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)
    if p.runs:
        p.runs[0].text = text
    else:
        p.add_run().text = text


def write_body(slide, lines):
    """lines: [(레벨, 글자)] — 0='1)' 소제목, 1='① 항목', 2='→ 부연'. 목차 장에서만 쓴다."""
    tf = body_of(slide).text_frame
    paras = tf.paragraphs
    proto = [copy.deepcopy(paras[0]._p), copy.deepcopy(paras[1]._p), copy.deepcopy(paras[2]._p)]
    txBody = paras[0]._p.getparent()
    for p in list(txBody.findall(NS + "p")):
        txBody.remove(p)
    for lvl, text in lines:
        p = copy.deepcopy(proto[lvl])
        runs = p.findall(NS + "r")
        for r in runs[1:]:
            p.remove(r)
        if not runs:
            p.append(copy.deepcopy(proto[lvl].findall(NS + "r")[0]))
            runs = p.findall(NS + "r")
        t = runs[0].find(NS + "t")
        t.text = text
        t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        txBody.append(p)


def picture(slide, name, x, y, w, cap=None, ratio=None):
    img = A(name)
    h = w / ratio
    slide.shapes.add_picture(img, Inches(x), Inches(y), Inches(w), Inches(h))
    if cap:
        D.text(slide, x, y + h + 0.1, w, 0.3, [(cap, {"size": 10, "color": D.INK2})],
               align=PP_ALIGN.CENTER)
    return y + h


# ── 1. 표지 ──────────────────────────────────────────────────
s = prs.slides[0]
for sh in s.shapes:
    if sh.has_text_frame and "LFP" in sh.text_frame.text:
        set_text(sh, "SafeAlert  근접 경보 앱")
    elif sh.has_text_frame and "2026.08" in sh.text_frame.text:
        p = sh.text_frame.paragraphs[0]
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
        p.runs[0].text = "2026.08 |  Flex Fulfillment Dynamic Ops  Waterflex"

# ── 2. Content ──────────────────────────────────────────────
write_body(prs.slides[1], [
    (0, "한 장 요약"),
    (0, " "),
    (0, "1. 개발 배경"),
    (0, "2. 기대효과"),
    (1, "  1) 사각지대 안전사고 사전 예방     2) 비용 비교     3) 투자 회수 (ROI)"),
    (0, "3. 장단점"),
    (1, "  1) I-PAS 와 비교     2) 앱의 단점     3) 대응 — PDA 이식"),
    (0, "4. 앞으로 수정계획"),
    (1, "  Phase 1 ~ 5"),
    (0, "5. 테스트 중간 결과"),
    (1, "  실 테스트 센터 (WF11 · WF21 · WF25)"),
    (0, " "),
    (0, "부록"),
    (1, "  A 비용 산출 근거     B 사용자 피드백     C 시연     D Andon     E 요청 사항"),
])

# ── 3. 한 장 요약 ───────────────────────────────────────────
s = prs.slides[2]
set_text(title_of(s), "한 장 요약")
D.clear_body(body_of(s))
D.text(s, 0.7, 1.12, 11.93, 0.6,
       [("사각지대를 사람의 주의력에 맡기지 않는다", {"bold": True, "size": 27})])
D.text(s, 0.7, 1.80, 11.93, 0.36,
       [("I-PAS 가 덮지 못하는 구간 — 보행자끼리 · EPJ · 미장착 장비 — 을 앱이 메운다", {"size": 13, "color": D.INK2})])
FACTS = [
    ("3개월 · 70회+", "현장에 배포하며 다듬은 기간"),
    ("센터당 331만원", "같은 범위를 하드웨어로 덮을 때"),
    ("테스트 3개 센터", "WF11 · WF21 · WF25 서베이 진행"),
    ("로드맵 1 / 5", "신뢰성 확보 단계 · 전 단계 출하 가능"),
]
for i, (big, cap) in enumerate(FACTS):
    x = 0.85 + i * 3.0
    D.text(s, x, 2.42, 2.9, 0.42, [(big, {"bold": True, "size": 18, "color": D.RED})])
    D.text(s, x, 2.88, 2.9, 0.34, [(cap, {"size": 11, "color": D.INK2})])

box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.62), Inches(11.93), Inches(2.62))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xF9)
box.line.color.rgb = D.LINE; box.line.width = Pt(1); box.shadow.inherit = False
box.text_frame.text = ""
D.text(s, 1.05, 3.88, 11.2, 0.42,
       [("요청  —  ", {"bold": True, "size": 15, "color": D.RED}),
        ("SafeAlert 를 정식 과제로 등록하고, 유지 · 검증에 필요한 리소스를 배정해 주십시오", {"bold": True, "size": 15})])
for i, t in enumerate([
    "①  유지보수 주체 지정 — 지금은 개인 작업물이라 승계 경로가 없다",
    "②  PDA 이식 — 개인 스마트폰 임시 설치를 업무 단말로 옮긴다 (Phase 2)",
    "③  현장 검증 슬롯 — 실기 검증이 유일한 회귀 확인 수단, 단계당 확인 1건",
    "④  보안 검토 — 릴리스 난독화 미적용 · 경보 이력 평문 저장 (확산 전 선행)",
]):
    D.text(s, 1.15, 4.45 + i * 0.38, 11.0, 0.36, [(t, {"size": 12})])
D.text(s, 1.05, 6.42, 11.2, 0.34,
       [("결정이 필요한 시점 — Phase 2 착수 전.  그때까지는 현재 버전이 현장에서 그대로 운영된다.",
         {"size": 11.5, "color": D.INK2, "italic": True})])

# ── 4. 1. 개발 배경 ─────────────────────────────────────────
s = prs.slides[3]
set_text(title_of(s), "1. 개발 배경")
D.clear_body(body_of(s))
picture(s, "fig_aisle.jpg", 8.62, 1.12, 4.28, "3m 렉 사이 통로 — 교차 구간에서 서로 보이지 않는다", 1200 / 750)
D.pill(s, 0.7, 1.12, "문제의 배경")
D.kv(s, 0.95, 1.64, 7.5, [
    ("AS-IS", "3m 높이 철제 렉 사이 통로에서 지게차와 보행자가 서로를 볼 수 없음"),
    ("Pain Point", "적재 파렛트가 2.4GHz 전파를 흡수하고, 장비 소음으로 경적 인지도 늦음"),
])
D.pill(s, 0.7, 2.52, "현재 상황")
D.kv(s, 0.95, 3.04, 7.5, [
    ("Data", "PIT 에는 I-PAS 가 장착돼 3m 에서 자동 셧다운한다"),
    ("Issue", "보행자끼리 · EPJ · 미장착 구간은 비어 있고, 그 구간은 주의력에만 맡겨져 있다"),
])
D.pill(s, 0.7, 4.02, "왜 앱으로 만들었나", w=2.05)
for i, t in enumerate([
    "①  안전사고 예방을 위해 I-PAS 와 같은 기능이 비어 있는 구간에도 필요하다고 판단",
    "②  장비마다 태그를 다는 방식으로는 사람과 구역을 다 덮지 못한다 — 지급된 단말끼리 직접 감지",
    "③  v1.0.1 부터 3개월간 70회 이상 배포하며 현장에서 계속 다듬음",
]):
    D.text(s, 0.95, 4.54 + i * 0.35, 11.9, 0.34, [(t, {})], size=12)
D.headline(s, 5.72, ["사각지대 해소", "인프라 없는 확산", "오경보 억제", "조용한 실패 제거"])
D.quote(s, 6.32, "I-PAS 가 장비를 지키는 동안, 그 밖의 구간은 아직 사람의 주의력에 기대고 있습니다.")

# ── 5. 2. 기대효과 (1) 사각지대 예방 ────────────────────────
s = prs.slides[4]
set_text(title_of(s), "2. 기대효과")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "1) 사각지대 예방", w=2.0)
D.text(s, 2.85, 1.15, 9.9, 0.3,
       [("I-PAS 와 비슷한 기능을 구현해, 사각지대에서 생기는 안전사고를 사전에 예방한다.",
         {"size": 12, "color": D.INK2})])
D.table(s, 0.72, 1.60, 12.1,
        ["구분", "As-is", "To-be"],
        [["감지", "사람의 주의력에만 의존", "소리 · 진동 · 화면 3중 경보로 보완"],
         ["도입", "안전 설비를 깔려면 장비별 장착 작업", "앱 설치만으로 즉시 적용"],
         ["확산", "장비를 늘려도 감지 범위는 그대로", "단말이 늘수록 감지쌍이 배로 (3대 3쌍 → 10대 45쌍)"],
         ["센터 차이", "센터마다 다른 조건에 대응 못 함", "판정 반경 · 신호 임계를 설정으로 조정"],
         ["후진 · 하역", "접근을 육안으로만 확인", "거리와 무관하게 즉시 최고 등급으로 경보"],
         ["경보 피로", "잦으면 작업자가 꺼버림", "세이프존 · 자동 음소거로 울릴 곳에서만"]],
        col_w=[1.55, 4.85, 5.70], row_h=0.55, size=11.5)
picture(s, "fig_pairs.png", 4.66, 5.32, 4.0, None, 728 / 278)

# ── 6. 2. 기대효과 (2) 비용 비교 ────────────────────────────
s = prs.slides[5]
set_text(title_of(s), "2. 기대효과")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "2) 비용 비교", w=1.7)
D.text(s, 2.55, 1.15, 10.2, 0.3,
       [("같은 범위를 하드웨어로 덮을 때 드는 금액이다. 실제 절감액이 아니다.",
         {"size": 12, "color": D.INK2})])
D.table(s, 0.72, 1.60, 12.1,
        ["항목", "단가", "일반 WF 센터 수량", "금액"],
        [["차량 태그", "525,000원", "PIT 4대", "2,100,000원"],
         ["EOD (3m 자동 셧다운)", "220,000원", "PIT 4대 전량", "880,000원"],
         ["보행자 태그", "110,000원", "상시 3명", "330,000원"],
         ["합계 — 센터 1곳", "", "", "3,310,000원"]],
        col_w=[3.30, 2.80, 3.00, 3.00], row_h=0.52, size=11.5,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT])
D.text(s, 0.95, 4.16, 11.9, 0.32,
       [("상시 인원 5명 센터는 ", {"size": 12}), ("3,530,000원", {"bold": True, "size": 12, "color": D.RED})])

D.pill(s, 0.7, 4.44, "센터 1곳 3,310,000원 구성 (만원)", w=3.75)
D.chart(s, 0.62, 4.80, 12.2, 1.12,
        ["차량 태그  52.5만 × 4대", "EOD  22만 × 4대", "보행자 태그  11만 × 3명"],
        [210, 88, 33], kind="bar", label_size=11, cat_size=10.5, headroom=1.12)
for i, t in enumerate([
    "·  자동 셧다운은 SafeAlert 가 못 한다. I-PAS 를 걷어내자는 이야기가 아니다.  단가는 사용자 제공 기준 (2026.08).",
    "·  개인 스마트폰 임시 설치 중 — PDA 로 옮기면 단말 비용 없이 업무 단말에서 돈다.  대체 여부는 로드맵 완료 후 재논의한다.",
]):
    D.text(s, 0.95, 6.00 + i * 0.30, 11.6, 0.29, [(t, {"size": 11, "color": D.INK2})])


# ── 7. 2. 기대효과 (3) 투자 회수 (ROI) ──────────────────────
s = prs.slides[6]
set_text(title_of(s), "2. 기대효과")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "3) 투자 회수 (ROI)", w=2.2)
D.text(s, 3.05, 1.15, 9.7, 0.3,
       [("절감액이 아니다. 하드웨어를 사야 덮이는 범위를 추가 구매 없이 덮는 데서 오는 회피 금액이다.",
         {"size": 12, "color": D.INK2})])
D.table(s, 0.72, 1.60, 6.45,
        ["도입 범위", "대상", "하드웨어 환산"],
        [["센터 1곳", "PIT 4대 · 상시 3명", "3,310,000원"],
         ["5개 센터", "PIT 20대 · 15명", "16,550,000원"],
         ["10개 센터", "PIT 40대 · 30명", "33,100,000원"],
         ["17개 센터 (전체)", "PIT 68대 · 51명", "56,270,000원"]],
        col_w=[2.05, 2.30, 2.10], row_h=0.48, size=11, head_size=11,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT])
D.chart(s, 7.30, 1.50, 5.55, 2.46,
        ["1곳", "5개", "10개", "17개"], [331, 1655, 3310, 5627],
        label_size=11, cat_size=11)
D.text(s, 7.30, 3.92, 5.55, 0.28,
       [("도입 센터 수에 따른 하드웨어 환산 누계 (만원)", {"size": 10, "color": D.INK3})],
       align=PP_ALIGN.CENTER)

D.pill(s, 0.7, 4.16, "투자 · 회수", w=1.55)
ROI = [
    ("투자 금액", "0원", "추가 하드웨어 없음 — 기존 단말에 앱을 설치한다"),
    ("회수 시점", "도입 즉시", "구매가 없으므로 회수 기간이 발생하지 않는다"),
    ("미산정", "공수 · 보완 태그", "개발 · PDA 이식 공수는 내부 인력, 금액 환산 전"),
]
for i, (k, v, sub) in enumerate(ROI):
    x = 0.72 + i * 4.09
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.56), Inches(3.92), Inches(0.98))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xF9)
    box.line.color.rgb = D.LINE; box.line.width = Pt(1); box.shadow.inherit = False
    box.text_frame.text = ""
    D.text(s, x + 0.18, 4.68, 3.6, 0.32,
           [(k + "   ", {"size": 10.5, "color": D.INK2}),
            (v, {"bold": True, "size": 15, "color": D.RED})])
    D.text(s, x + 0.18, 5.08, 3.6, 0.40, [(sub, {"size": 10, "color": D.INK2})])
for i, t in enumerate([
    "·  단가는 사용자 제공 기준 (2026.08) — 차량 태그 52.5만 · EOD 22만 · 보행자 태그 11만.",
    "·  EOD 의 3m 자동 셧다운은 SafeAlert 가 대신하지 못한다. 금액 비교는 감지 · 경보 범위에 한한다.",
    "·  iOS 인원 보완 태그(11만원/개)와 세이프존 비콘은 대상 인원 확정 전이라 산입하지 않았다.",
]):
    D.text(s, 0.95, 5.68 + i * 0.31, 11.5, 0.30, [(t, {"size": 11, "color": D.INK2})])

# ── 8. 3. 장단점 ────────────────────────────────────────────
s = prs.slides[7]
set_text(title_of(s), "3. 장단점")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "1) I-PAS 와 비교", w=2.0)
D.table(s, 0.72, 1.58, 12.1,
        ["구분", "I-PAS", "SafeAlert"],
        [["동작", "3m 내 감지 시 자동 셧다운 — 최후 개입", "15m 경고 · 8m 위험 경보. 장비 정지는 못 한다"],
         ["덮는 대상", "태그를 단 장비와 사람", "앱을 넣은 모든 단말 (보행자끼리 · EPJ 포함)"],
         ["대상 확대", "사람 · 장비가 늘면 태그를 그만큼 더 산다", "앱 설치만 — 추가 구매 없음"],
         ["유지보수", "하드웨어 고장 · 배터리 교체", "앱 업데이트 (자동 배포)"]],
        col_w=[1.70, 5.20, 5.20], row_h=0.50, size=11.5)

D.pill(s, 0.7, 4.06, "2) 앱의 단점", w=1.75)
for i, t in enumerate([
    "①  지게차 셧다운 불가 — 경보만 한다. I-PAS EOD 의 3m 자동 정지는 대체하지 못한다",
    "②  일시적 개인 스마트폰 설치 — 업무 단말이 아니라 개인 폰에 깔려 있어 지속 가능한 상태가 아니다",
    "③  비콘 구매 필요 — 세이프존(무음구역)과 미설치 인원 감지에는 BLE 비콘 · 태그가 든다",
    "④  안드로이드 전용 · 20대 이상 밀집 시 화면 반응 저하",
]):
    D.text(s, 0.95, 4.54 + i * 0.34, 11.9, 0.32, [(t, {"size": 11.5})])

box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(6.02), Inches(12.1), Inches(0.76))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF8, 0xF9)
box.line.color.rgb = D.LINE; box.line.width = Pt(1); box.shadow.inherit = False
box.text_frame.text = ""
D.text(s, 1.0, 6.18, 11.6, 0.5,
       [("3) 대응  —  ", {"bold": True, "size": 13, "color": D.RED}),
        ("Phase 2 에 PDA 로 이식하는 것이 가장 효과적이다. ", {"bold": True, "size": 13}),
        ("개인 폰 의존과 배터리 부담이 함께 없어지고, 업무 단말이라 배포 · 회수 경로도 이미 있다.",
         {"size": 12, "color": D.INK2})])

# ── 9. 4. 앞으로 수정계획 ───────────────────────────────────
s = prs.slides[8]
set_text(title_of(s), "4. 앞으로 수정계획")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "Phase 1 ~ 5", w=1.6)
D.text(s, 2.45, 1.15, 10.2, 0.3,
       [("새 기능을 붙이는 작업이 아니다. 판정이 흔들리는 구조적 원인을 걷어낸다.", {"size": 12, "color": D.INK2})])
D.table(s, 0.72, 1.62, 12.1,
        ["단계", "무엇을 하는가", "출하 시 현장에서 확인할 것", "상태"],
        [["1", "테스트 하네스와 CI 회귀 게이트", "설치 · 경보 동작이 이전과 같은가", "완료"],
         ["2", "안전 크리티컬 경로 골든 테스트 + PDA 이식", "천천히 접근하는 지게차에 경고가 뜨는가", "예정"],
         ["3", "BleService 분해", "달라진 게 없는가 (동작 보존)", "예정"],
         ["4", "기기 상태 단일화", "2시간 이상 구동 후에도 느려지지 않는가", "예정"],
         ["5", "판정 워커 분리", "20대 이상에서 화면이 끊기지 않는가", "예정"]],
        col_w=[0.75, 4.65, 5.40, 1.30], row_h=0.52, size=11.5,
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER])
D.pill(s, 0.7, 4.72, "이 순서인 이유", w=1.6)
for i, t in enumerate([
    "①  측정 수단이 먼저다 — 테스트 0건 상태에서 안전 로직을 분해하면, 회귀를 현장에서 발견하는 방식이 그대로 남는다",
    "②  단계마다 단독 배포한다 — 변경을 섞어 내보내면 현장에서 회귀가 나도 어느 변경 탓인지 특정할 수 없다",
]):
    D.text(s, 0.95, 5.22 + i * 0.42, 11.9, 0.40, [(t, {"size": 11.5})])
D.headline(s, 6.22, ["완료 1 / 5", "전 단계 출하 가능", "PDA 이식은 Phase 2"])

# ── 10. 5. 테스트 중간 결과 ──────────────────────────────────
s = prs.slides[9]
set_text(title_of(s), "5. 테스트 중간 결과")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "실 테스트 센터", w=1.8)
D.text(s, 2.65, 1.15, 10.1, 0.3,
       [("WF11 · WF21 · WF25 에서 실사용 테스트 진행 중", {"size": 12, "color": D.INK2})])
D.table(s, 0.72, 1.60, 12.1,
        ["서베이 항목 (안)", "WF11", "WF21", "WF25"],
        [["경보 타이밍 — 너무 이르거나 늦지 않은가", "—", "—", "—"],
         ["거리 정확도 — 표시 거리가 체감과 맞는가", "—", "—", "—"],
         ["오경보 빈도 — 울릴 필요 없는데 울리는가", "—", "—", "—"],
         ["배터리 · 발열 — 1교대를 버티는가", "—", "—", "—"],
         ["사용 편의 — 역할 전환 · 알림 조작", "—", "—", "—"],
         ["종합 만족도", "—", "—", "—"]],
        col_w=[6.10, 2.00, 2.00, 2.00], row_h=0.50, size=11.5,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER])

box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(5.06), Inches(12.1), Inches(0.80))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFD, 0xF3, 0xF2)
box.line.color.rgb = D.RED; box.line.width = Pt(1); box.shadow.inherit = False
box.text_frame.text = ""
D.text(s, 1.0, 5.24, 11.6, 0.5,
       [("작성 전  —  ", {"bold": True, "size": 13, "color": D.RED}),
        ("서베이를 회수 · 집계한 뒤 이 장을 채운다. 항목은 제안이며 확정본이 아니다. "
         "채우기 전에는 이 장을 빼고 보고한다.", {"size": 12})])
for i, t in enumerate([
    "·  5점 척도(5 매우 만족 ~ 1 매우 미흡)로 집계하면 LFP 지게차 평가 자료와 형식이 맞는다.",
    "·  센터별 응답자 수와 테스트 기간을 함께 적어야 표본 크기를 판단할 수 있다.",
]):
    D.text(s, 0.95, 6.10 + i * 0.34, 11.5, 0.32, [(t, {"size": 11.5, "color": D.INK2})])

# ── 11. 부록 A — 비용 산출 근거 ─────────────────────────────
s = prs.slides[10]
set_text(title_of(s), "부록 A — 비용 산출 근거")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "확산 환산", w=1.4)
D.text(s, 2.3, 1.15, 10.4, 0.3,
       [("센터 1곳 3,310,000원 (일반 WF · PIT 4대 · 상시 3명) 기준", {"size": 12, "color": D.INK2})])
D.table(s, 0.72, 1.60, 12.1,
        ["센터 수", "5개", "10개", "17개 (FFWF 전체)"],
        [["같은 범위를 하드웨어로 덮을 때", "16,550,000원", "33,100,000원", "56,270,000원"]],
        col_w=[4.30, 2.60, 2.60, 2.60], row_h=0.52, size=11.5,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
for i, t in enumerate([
    "·  단가는 사용자 제공 (2026.08 기준). 계약 조건에 따라 달라질 수 있다.",
    "·  SafeAlert 는 같은 대상(PIT 4대 + 상시 인원)을 단말에 앱을 넣어 덮는다. 추가 하드웨어 구매 없음.",
    "·  다만 EOD 의 3m 자동 셧다운은 SafeAlert 가 대신하지 못한다 — 금액 비교는 감지 · 경보 범위에 한한다.",
    "·  세이프존(무음구역)과 미설치 인원 감지에는 BLE 비콘 · 태그 구매가 별도로 든다.",
]):
    D.text(s, 0.95, 2.66 + i * 0.34, 11.5, 0.32, [(t, {"size": 11.5, "color": D.INK2})])
D.pill(s, 0.72, 4.06, "센터 1곳 구성 (만원)", w=2.5)
D.chart(s, 0.62, 4.44, 6.5, 1.86,
        ["차량 태그  52.5만 × 4", "EOD  22만 × 4", "보행자 태그  11만 × 3"],
        [210, 88, 33], kind="bar", label_size=11, cat_size=10, headroom=1.12)
picture(s, "fig_radius.png", 7.90, 4.14, 3.5, "역할 조합별 판정 반경 (실제 비례 축척)", 648 / 358)

# ── 12. 부록 B — 사용자 피드백 ──────────────────────────────
s = prs.slides[11]
set_text(title_of(s), "부록 B — 사용자 피드백")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "오더십 리뷰 반영", w=2.0)
D.table(s, 0.72, 1.60, 8.3,
        ["차수", "반영 내용", "결과"],
        [["1차", "역할 전환 버튼 추가 · EPJ 역할 버튼 숨김 · 스플래시 교체", "앱 재설치 없이 화면에서 역할 전환"],
         ["2차", "같은 기기 · 같은 등급에 5초 이상 머물면 소리 · 진동 자동 음소거", "표시 · 기록은 유지, 등급 오르면 즉시 재발령"],
         ["3차", "세이프존 — 휴게실 · 충전 구역에 비콘을 두면 경보 억제", "존 안의 기기는 상대에게도 안전으로 인식"]],
        col_w=[0.85, 4.35, 3.10], row_h=0.62, size=11)
picture(s, "fig_zone.png", 9.28, 1.52, 3.48, None, 668 / 288)
D.pill(s, 0.7, 4.02, "현장 요구 반영", w=1.85)
D.table(s, 0.72, 4.50, 12.1,
        ["항목", "Before", "After"],
        [["보호 상태", "끊겨도 작업자가 알 수 없음", "감지가 멈추면 상시 알림에 이상 표시"],
         ["이탈 후", "멀어졌는데 다시 울림", "경고 범위 재발령 억제"],
         ["알림 끄기", "단계가 번거로움", "알림 본문을 누르면 즉시 무음"],
         ["경보 화면", "눈에 잘 들어오지 않음", "화면 가장자리에 붙는 사이드바"]],
        col_w=[1.75, 4.65, 5.70], row_h=0.44, size=11.5)

# ── 13. 부록 C — 시연 ───────────────────────────────────────
s = prs.slides[12]
set_text(title_of(s), "부록 C — 시연")
D.clear_body(body_of(s))
D.flow(s, 0.7, 1.02, 12.2,
       ["신호 수신", "거리 추정", "등급 판정", "3중 경보"],
       ["BLE 광고", "3단 필터", "상태머신", "소리 · 진동 · 화면"])
ph_h = 3.80
ph_w = ph_h * 300 / 624
for k, (png, cap) in enumerate([("screen_running", "동작 화면"), ("screen_settings", "설정 화면")]):
    x = 1.15 + k * (ph_w + 0.5)
    D.text(s, x - 0.25, 2.12, ph_w + 0.5, 0.3, [(cap, {"bold": True, "size": 12})], align=PP_ALIGN.CENTER)
    s.shapes.add_picture(f"{REPO}/assets/{png}.png", Inches(x), Inches(2.48), Inches(ph_w), Inches(ph_h))
vw, vx, vy = 6.35, 6.42, 2.32
vh = vw * 868 / 1544
D.text(s, vx, vy - 0.38, vw, 0.32,
       [("시연 영상 — 슬라이드쇼에서 재생 (경보음 포함)", {"bold": True, "size": 12})], align=PP_ALIGN.CENTER)
s.shapes.add_movie(f"{REPO}/web/safealert-sim.mp4", Inches(vx), Inches(vy), Inches(vw), Inches(vh),
                   poster_frame_image=f"{REPO}/web/sim-poster.jpg", mime_type="video/mp4")
D.text(s, vx, vy + vh + 0.18, vw, 0.32,
       [("판정 시뮬레이터 — web/safealert-simulator.html 에서 직접 조작", {"size": 11, "color": D.INK2})],
       align=PP_ALIGN.CENTER)

# ── 14. 부록 D — Andon ──────────────────────────────────────
s = prs.slides[13]
set_text(title_of(s), "부록 D — Andon (이럴 때 멈춘다)")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "현장 운영 중", w=1.65)
D.pill(s, 6.95, 1.12, "업데이트 · 수정 중", w=2.15)
FIELD = [
    ("미발령", "접근했는데 경보가 뜨지 않은 사례 1건이라도", True),
    ("오발령 지속", "정지 중 계속 울림 · 이탈 후 다시 울림이 반복", True),
    ("감지 중단", "상시 알림이 ‘이상’ 으로 바뀌거나 사라짐", False),
    ("경보 피로", "무음으로 두거나 앱을 꺼 둔 단말이 늘어남", False),
]
DEPLOY = [
    ("버전 혼재", "구버전 단말이 신버전을 감지하지 못함", True),
    ("보호 공백", "작업 시간 중 일괄 업데이트 금지 — 교대 전환 때만", False),
    ("서비스 미기동", "업데이트 · 역할 전환 후 백그라운드 실행 미복구", False),
    ("회귀 재발", "이전에 고친 증상이 다시 관측됨", True),
]
for i, ((k1, v1, u1), (k2, v2, u2)) in enumerate(zip(FIELD, DEPLOY)):
    y = 1.62 + i * 0.84
    for x, w, (k, v, urgent) in ((0.9, 5.5, (k1, v1, u1)), (7.15, 5.6, (k2, v2, u2))):
        D.text(s, x, y, w, 0.78,
               [[(k, {"bold": True, "size": 12.5, "color": D.RED if urgent else D.INK}),
                 ("   즉시 중단" if urgent else "   관측 시 판단", {"size": 10, "color": D.INK3})],
                [(v, {"size": 11.5, "color": D.INK2})]], space=2)
D.pill(s, 0.7, 4.98, "공통 조치", w=1.4)
for i, (k, v) in enumerate([
    ("즉시 롤백", "직전 태그 APK 로 되돌린다 — 현장 재설치 없이 배포된다"),
    ("선행 확인", "전 단말 배포 전 1대에서 먼저 확인한다"),
    ("빌드 차단", "골든 테스트가 실패하면 릴리스가 자동 차단된다 (동작 중)"),
    ("원인 귀속", "단계를 단독 배포하므로 어느 변경 탓인지 특정된다"),
    ("재개 조건", "실패한 상황을 테스트로 고정한 뒤에만 다음 단계로"),
]):
    col, row = divmod(i, 3)
    D.text(s, 0.95 + col * 6.25, 5.44 + row * 0.35, 5.9, 0.33,
           [(k + "  —  ", {"bold": True, "size": 11.5}), (v, {"size": 11.5, "color": D.INK2})])
D.quote(s, 6.56, "멈출 기준을 먼저 정해 두는 것이, 멈추지 않고 밀어붙이는 것보다 빠릅니다.", x=0.7, w=8.4)

# ── 15. 부록 E — 요청 사항 ──────────────────────────────────
s = prs.slides[14]
set_text(title_of(s), "부록 E — 요청 사항")
D.clear_body(body_of(s))
D.pill(s, 0.7, 1.12, "요청", w=1.1)
D.text(s, 2.0, 1.15, 10.8, 0.36,
       [("SafeAlert 를 정식 과제로 등록하고, 유지 · 검증에 필요한 리소스를 배정해 주십시오.", {"bold": True, "size": 14})])
D.table(s, 0.72, 1.72, 12.1,
        ["무엇을", "왜 필요한가", "없으면"],
        [["정식 과제 등록", "지금은 개인 작업물이다. 유지보수 주체가 지정돼 있지 않다", "만든 사람이 빠지면 멈춘다"],
         ["PDA 이식", "개인 스마트폰 임시 설치를 업무 단말로 옮겨야 한다 (Phase 2)", "개인 폰 의존이 계속된다"],
         ["현장 검증 슬롯", "실기 검증이 유일한 회귀 확인 수단 — 단계당 확인 항목 1건", "검증 없이 배포하거나 배포가 멈춘다"],
         ["보안 검토", "릴리스 빌드 난독화 미적용 · 경보 이력 평문 저장", "확산할수록 리스크가 커진다"]],
        col_w=[2.30, 5.70, 4.10], row_h=0.60, size=11.5)
D.pill(s, 0.7, 4.92, "지금 상태", w=1.4)
for i, t in enumerate([
    "①  v1.1.70 이 현장에서 운영 중이다 — 결정이 늦어도 보호가 멈추지는 않는다",
    "②  다만 신뢰성 로드맵은 1 / 5 에서 멈춰 있고, 남은 네 단계가 재발 고리를 끊는 작업이다",
]):
    D.text(s, 0.95, 5.40 + i * 0.38, 11.9, 0.36, [(t, {"size": 12})])
D.quote(s, 6.34, "지금 필요한 것은 새 기능이 아니라, 이미 도는 것을 믿을 수 있게 만드는 시간입니다.", x=0.7, w=8.4)

prs.save(OUT)
print("wrote", OUT)
