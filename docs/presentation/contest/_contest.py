# CouSolve Idea Contest 제출 양식에 내용만 채우기 위한 도구.
# 폼 자체 — 상단 단계 탭 · 흰 둥근 카드 · 회색 라벨 칩 · 파란 장식 — 은 손대지 않는다.
# 글꼴(맑은 고딕)과 강조색(#053866)은 양식이 쓰던 것을 그대로 따른다.
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

NAVY = RGBColor(0x05, 0x38, 0x66)   # 양식의 활성 탭 색
INK = RGBColor(0x40, 0x40, 0x40)    # 양식 본문 색
INK2 = RGBColor(0x6B, 0x72, 0x78)
INK3 = RGBColor(0x9A, 0xA0, 0xA6)
LINE = RGBColor(0xD8, 0xDD, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
# 2계열 비교용 — dataviz validate_palette.js (light) 6/6 PASS
TEAL = RGBColor(0x1B, 0x8A, 0x6B)
AMBER = RGBColor(0xB9, 0x77, 0x0B)
FONT = "맑은 고딕"


def _run(p, t, size, bold=False, color=INK, italic=False, u=False, font=FONT):
    """u=True 는 추정치 표시다. 자료 전체에서 밑줄은 '실측 아님' 하나의 뜻만 갖는다."""
    r = p.add_run()
    r.text = t
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.underline = bool(u)
    r.font.name = font
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = parse_xml(f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            rPr.append(e)
        e.set("typeface", font)
    return r


def _norm(rows):
    """문자열 · chunk 리스트 · 문단 리스트를 모두 '문단 리스트' 로 맞춘다."""
    if isinstance(rows, str):
        return [[(rows, {})]]
    if rows and isinstance(rows[0], tuple):      # 한 문단짜리 chunk 리스트
        return [rows]
    return rows


def write(tf, rows, size=11, align=PP_ALIGN.LEFT, space=1.25, gap=0):
    """텍스트 프레임을 비우고 문단을 다시 쌓는다. 도형의 채우기 · 테두리는 건드리지 않는다."""
    rows = _norm(rows)
    body = tf.paragraphs[0]._p.getparent()
    NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for p in list(body.findall(NS + "p"))[1:]:
        body.remove(p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for i, row in enumerate(rows):
        p = p0 if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        if gap and i:
            p.space_before = Pt(gap)
        for t, o in row:
            _run(p, t, o.get("size", size), o.get("bold", False), o.get("color", INK),
                 o.get("italic", False), o.get("u", False))
    return tf


# 라벨 칩이 상자 왼쪽 위를 덮고 있다. 그 아래부터 글이 시작해야 한다.
TOPCHIP = 0.40


def fill(shape, rows, size=11, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         margin=0.18, space=1.25, gap=0, top=0.10):
    """양식에 이미 있는 상자에 글자만 넣는다. top 은 라벨 칩을 피하는 윗 여백."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = anchor
    return write(tf, rows, size, align, space, gap)


def body_of(slide, label):
    """라벨 칩 글자로 그룹을 찾아, 같은 그룹의 큰 상자를 돌려준다."""
    for sh in slide.shapes:
        if sh.shape_type != 6:
            continue
        subs = list(sh.shapes)
        hit = [s for s in subs if s.has_text_frame and s.text_frame.text.strip() == label]
        if hit:
            rest = [s for s in subs if s is not hit[0]]
            return max(rest, key=lambda s: (s.width or 0) * (s.height or 0))
    raise SystemExit(f"라벨을 찾지 못했다: {label!r}")


def named(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise SystemExit(f"도형을 찾지 못했다: {name!r}")


def at(slide, x, y, tol=0.05):
    """좌표로 도형을 집는다 (그룹이 아닌 낱개 도형용)."""
    for sh in slide.shapes:
        if sh.left is None:
            continue
        if abs(sh.left / 914400 - x) < tol and abs(sh.top / 914400 - y) < tol:
            return sh
    raise SystemExit(f"({x}, {y}) 위치의 도형을 찾지 못했다")


def text(slide, x, y, w, h, rows, size=11, align=PP_ALIGN.LEFT, space=1.2):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    write(tf, rows, size, align, space)
    return box


def lines(slide, x, y, w, items, size=10.5, gap=0.28, bullet=""):
    for i, it in enumerate(items):
        rows = it if not isinstance(it, str) else [(bullet + it, {})]
        text(slide, x, y + i * gap, w, gap, rows, size=size)
    return y + len(items) * gap


def _border(c, edge, color, pt=0.75):
    tag = {"L": "a:lnL", "R": "a:lnR", "T": "a:lnT", "B": "a:lnB"}[edge]
    tcPr = c._tc.get_or_add_tcPr()
    for e in tcPr.findall(qn(tag)):
        tcPr.remove(e)
    ln = parse_xml(
        f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'w="{int(pt * 12700)}" cap="flat" cmpd="sng" algn="ctr">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'<a:prstDash val="solid"/></{tag}>')
    order = ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]
    after = None
    for t in order[order.index(tag) + 1:]:
        after = tcPr.find(qn(t))
        if after is not None:
            break
    (after.addprevious(ln) if after is not None else tcPr.insert(0, ln))


def cell(c, rows, size=9.5, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
         color=INK, margin=0.06, space=1.15):
    tf = c.text_frame
    tf.word_wrap = True
    c.margin_left = c.margin_right = Inches(margin)
    c.margin_top = c.margin_bottom = Inches(0.02)
    c.vertical_anchor = anchor
    rows = _norm(rows)
    if color is not INK:
        rows = [[(t, {**o, "color": o.get("color", color)}) for t, o in r] for r in rows]
    return write(tf, rows, size, align, space)


def table(slide, x, y, w, header, rows, col_w, row_h=0.28, head_h=0.28,
          size=9.5, head_size=9.5, aligns=None, head_fill=NAVY):
    n_r, n_c = len(rows) + 1, len(header)
    gf = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w),
                                Inches(head_h + row_h * len(rows)))
    tbl = gf.table
    pr = tbl._tbl.find(qn("a:tblPr"))
    pr.set("bandRow", "0"); pr.set("firstRow", "1")
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, n_r):
        tbl.rows[i].height = Inches(row_h)
    for j, h in enumerate(header):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = head_fill
        cell(c, [[(h, {"bold": True, "color": WHITE})]], head_size, PP_ALIGN.CENTER, margin=0.05)
        for e in "LRTB":
            _border(c, e, "053866")
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE
            cell(c, v, size, (aligns or [PP_ALIGN.LEFT] * n_c)[j], margin=0.05)
            for e in "LRT":
                _border(c, e, "FFFFFF")
            _border(c, "B", "D8DDE1")
    return gf


def _quiet(ch, cat_size, label_size, num_fmt):
    from pptx.enum.chart import XL_LABEL_POSITION, XL_TICK_MARK
    ch.has_title = False
    ch.font.name = FONT; ch.font.size = Pt(cat_size); ch.font.color.rgb = INK2
    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = num_fmt; dl.number_format_is_linked = False
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.font.name = FONT; dl.font.size = Pt(label_size)
    dl.font.bold = True; dl.font.color.rgb = INK
    va = ch.value_axis
    va.visible = False
    va.has_major_gridlines = False; va.has_minor_gridlines = False
    va.minimum_scale = 0
    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.major_tick_mark = XL_TICK_MARK.NONE; ca.minor_tick_mark = XL_TICK_MARK.NONE
    ca.format.line.color.rgb = LINE
    ca.tick_labels.font.name = FONT; ca.tick_labels.font.size = Pt(cat_size)
    ca.tick_labels.font.color.rgb = INK2
    cs = ch._chartSpace
    C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
    cs.insert(list(cs).index(cs.find(C + "chart")) + 1, parse_xml(
        '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>'))
    return plot


def chart2(slide, x, y, w, h, cats, series, num_fmt="#,##0", label_size=9,
           cat_size=9, colors=(AMBER, TEAL), gap=80, headroom=1.22):
    """계열 두 개짜리 세로 막대. 계열이 둘이라 범례를 반드시 둔다."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    cd = CategoryChartData()
    cd.categories = list(cats)
    for name, vals in series:
        cd.add_series(name, tuple(vals), number_format=num_fmt)
    ch = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y),
                                Inches(w), Inches(h), cd).chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.TOP
    ch.legend.include_in_layout = False
    ch.legend.font.name = FONT; ch.legend.font.size = Pt(cat_size)
    ch.legend.font.color.rgb = INK2
    plot = _quiet(ch, cat_size, label_size, num_fmt)
    plot.gap_width = gap; plot.overlap = -12
    for ser, col in zip(plot.series, colors):
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
        ser.format.line.fill.background()
    ch.value_axis.maximum_scale = max(max(v) for _, v in series) * headroom
    return ch


def kpi(slide, x, y, w, h, tag, name, desc, val, color, est=True):
    """색 머리띠 + 흰 본문. 값에 밑줄이면 추정치."""
    from pptx.enum.shapes import MSO_SHAPE
    head = 0.34
    top = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    top.fill.solid(); top.fill.fore_color.rgb = WHITE
    top.line.color.rgb = LINE; top.line.width = Pt(0.75); top.shadow.inherit = False
    top.text_frame.text = ""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.10), Inches(y + 0.11),
                                  Inches(0.055), Inches(h - 0.22))
    band.fill.solid(); band.fill.fore_color.rgb = color
    band.line.fill.background(); band.shadow.inherit = False
    band.text_frame.text = ""
    text(slide, x + 0.26, y + 0.07, w - 0.40, 0.20,
         [[(tag + "   ", {"size": 8, "color": INK3}),
           (name, {"bold": True, "size": 10.5, "color": color})]])
    text(slide, x + 0.26, y + 0.28, w - 0.40, 0.20, [[(desc, {"size": 8.5, "color": INK2})]])
    text(slide, x + 0.26, y + h - 0.30, w - 0.40, 0.24,
         [[(val, {"bold": True, "size": 12, "color": color, "u": est})]])
    return top
