# CouSolve Idea Contest 제출 양식(16장)에 SafeAlert 내용을 채운다.
#   python3 build_contest_deck.py <CouSolve_Contest_양식.pptx>
#
# 양식 원본은 사내 문서라 저장소에 넣지 않는다 — 로컬 경로를 인자로 넘긴다.
# 양식의 3 · 4장은 도형이 하나도 없는 빈 장(상단 탭도 없다)이라 걷어낸다. 남는 14장:
#   1 표지 / 2 01.결론 / 3 01.Summary / 4 Summary 자유 / 5 02.Problem
#   6 02.5Why / 7 03.Tradeoffs / 8 04.Benchmark / 9 05.Solution
#   10 Solution 표 / 11 Solution 자유 / 12 06.Metrics / 13 07.Andon / 14 08.Feedback
#
# 숫자 원칙 — 실측이 없는 칸은 비우지 않고 추정치를 넣되 **밑줄**로 표시한다.
# 자료 전체에서 밑줄은 '실측 아님' 한 가지 뜻만 갖는다.
# 경보 건수는 2026-09-02 부터 실측값이다 (ALERTS.md · 주간 자동 집계). 12장에
# 처음 추정과 나란히 두어, 하나뿐이던 가정이 얼마나 빗나갔는지 그대로 보인다.
# 아직 밑줄이 남은 칸: 쓸데없이 울린 횟수 · 작업자 만족도 (설문 회수 전).
import os, shutil, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import _contest as C

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(os.path.dirname(HERE), "company_form", "assets")
OUT = os.path.join(HERE, "SafeAlert_CouSolve_Contest.pptx")

# ── 추정치 ──────────────────────────────────────────────────
PAIRS, SHIFT_H, RATE = 21, 9, 0.5      # 감지 쌍 · 1교대 · 쌍당 시간당 조우(유일한 가정값)
_h = lambda x: int(x + 0.5)
WARN = _h(PAIRS * SHIFT_H * RATE)      # 21 × 9 × 0.5 = 94.5 → 95
RATIO = round((8 / 15) ** 2, 2)        # 반경 비의 제곱 0.284 → 0.28
DANGER = _h(WARN * RATIO)              # 95 × 0.28 = 26.6 → 27
LATEST, LATEST_DATE = "v1.1.71", "2026.09.01"   # github 릴리스 기준
RELEASES = 97          # git tag 실측 (v1.0.20 2026-06-07 ~ v1.1.71 2026-09-01)
EST = {"u": True}
ESTB = {"u": True, "bold": True}

# ── 실측치 ──────────────────────────────────────────────────
# .github/workflows/alert-digest.yml 이 매주 갱신하는 저장소 맨 위 ALERTS.md 에서 옮겼다.
# WF11 사업장 노드, 2026-08-04 ~ 09-01 중 경보 기록이 남은 12일.
# 1건 = 경보 1회가 아니라 '가까워진 1분'이다 (같은 상대는 1분에 한 번만 기록된다).
M_FROM, M_TO, M_DAYS = "2026.08.04", "2026.09.01", 12
M_WARN_T, M_DANGER_T = 487, 291                    # 기간 합계
M_WARN   = round(M_WARN_T / M_DAYS, 1)             # 40.6 회/일
M_DANGER = round(M_DANGER_T / M_DAYS, 1)           # 24.2 회/일
M_DEVICES, M_PAIRS = 41, 41
M_RATE  = round(M_WARN / (M_PAIRS * SHIFT_H), 2)   # 0.11 회/시간/짝  (추정 0.5)
M_RATIO = round(M_DANGER / M_WARN, 2)              # 0.60             (추정 0.28)

if len(sys.argv) < 2:
    raise SystemExit("사용법: python3 build_contest_deck.py <CouSolve_Contest_양식.pptx>")
BASE = os.path.join(HERE, "base_contest.pptx")
shutil.copyfile(os.path.abspath(sys.argv[1]), BASE)
prs = Presentation(BASE)
assert len(prs.slides) == 16, f"양식이 16장이어야 한다 (지금 {len(prs.slides)}장)"

# 빈 3 · 4장 제거 — 상단 탭도 없는 백지라 제출본에 그대로 두면 결번처럼 보인다
lst = prs.slides._sldIdLst
for i in (3, 2):
    rId = lst[i].rId
    prs.part.drop_rel(rId)
    del lst[i]
S = prs.slides


def pic(slide, name, x, y, w, ratio, cap=None):
    h = w / ratio
    slide.shapes.add_picture(os.path.join(ASSETS, name), Inches(x), Inches(y),
                             Inches(w), Inches(h))
    if cap:
        C.text(slide, x, y + h + 0.05, w, 0.22, [[(cap, {"size": 8, "color": C.INK3})]],
               align=PP_ALIGN.CENTER)
    return y + h


# ── 1. 표지 ─────────────────────────────────────────────────
s = S[0]
C.fill(C.named(s, "TextBox 1"),
       [[("지게차와 보행자가 서로 다가오면 미리 알려주는 앱",
          {"bold": True, "size": 20, "color": C.NAVY})],
        [("— 이미 나눠준 업무용 스마트폰끼리 신호를 주고받습니다 (SafeAlert)",
          {"size": 12, "color": C.INK2})]],
       margin=0.02, space=1.25)
C.fill(C.at(s, 1.74, 6.47), [[("WF11", {"bold": True, "size": 12})]],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.05)
C.fill(C.at(s, 5.84, 6.47), [[("Ian", {"bold": True, "size": 12})]],
       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.05)

# ── 2. 01. 결론 ─────────────────────────────────────────────
s = S[1]
CONC = [
    (2.11, "문제", "우리 센터에는 지게차 접근을 알리는 장치가 없습니다",
     [[("시야는 적재 파렛트에 막히고, 경적은 장비 소음에 묻힙니다. "
        "지게차 4대 · 상시 3명, 어느 쪽도 감지되지 않습니다. "
        "지금은 작업자의 주의력이 사실상 유일한 안전장치입니다.", {})]]),
    (3.63, "해결", "이미 나눠준 스마트폰끼리 서로를 감지하게 만들었습니다",
     [[("앱만 깔면 서로 다가올 때 소리 · 진동 · 화면으로 알립니다. "
        "같은 일을 장비로 하면 우리 센터 규모로 331만원, ", {}),
       ("이 앱은 사람이 늘어도 추가 비용이 없습니다", {"bold": True, "color": C.TEAL}),
       (". 지게차를 자동으로 멈추지는 못합니다.", {})]]),
    (5.19, "확인", "세 센터에서 지금 돌아가고 있는 앱입니다",
     [[("WF11 · WF21 · WF25 에서 실사용 중입니다. WF11 에서는 하루 평균 ", {}),
       (f"경고 {M_WARN}회 · 위험 {M_DANGER}회", {"bold": True, "color": C.NAVY}),
       (f"가 기록됩니다 (기록 {M_DAYS}일 평균 · 1건 = 서로 가까워진 1분).", {})],
      [("코드 · 배포 이력  ", {"size": 10, "color": C.INK3}),
       ("github.com/pslymzero-cyber/SafeAlert", {"bold": True, "color": C.NAVY, "size": 11.5}),
       ("      최신 버전 ", {"size": 10, "color": C.INK2}),
       (f"{LATEST}", {"bold": True, "size": 10.5, "color": C.TEAL}),
       (f"  ({LATEST_DATE} 배포)", {"size": 10, "color": C.INK2})]]),
]
for i, (y, tag, head, sub) in enumerate(CONC, start=1):
    C.fill(C.at(s, 1.85, y),
           [[(f"{tag}", {"bold": True, "size": 11, "color": C.INK3}),
             ("      ", {"size": 11}),
             (head, {"bold": True, "size": 15.5, "color": C.NAVY})]] + sub,
           size=11, anchor=MSO_ANCHOR.MIDDLE, margin=0.32, space=1.35, gap=5)

# ── 3. 01. Summary ──────────────────────────────────────────
s = S[2]
C.fill(C.body_of(s, "Background"),
       [[("우리 센터에는 지게차와 보행자가 서로 다가오는 것을 알려주는 장치가 없다. "
          "통로가 만나는 교차 구간에서는 적재 파렛트가 시야를 막고, 장비 소음에 경적이 묻힌다. "
          "서로를 보지 못한 채 가까워진다.", {"size": 10.5})]],
       margin=0.20, space=1.3, top=C.TOPCHIP)
C.fill(C.body_of(s, "Solution"),
       [[("이미 나눠준 스마트폰에 앱만 깔면 서로를 감지한다.",
          {"bold": True, "size": 11.5})],
        [("따로 살 장비도, 몸에 달 태그도, 설치할 것도 없다.",
          {"size": 10, "color": C.INK2})],
        [("신호 받기  →  거리 계산  →  위험도 판단  →  소리 · 진동 · 화면 경보",
          {"bold": True, "size": 10.5, "color": C.NAVY})],
        [("정밀 측위(UWB) 폰끼리는 실제 거리로 잰다. 지게차가 낀 짝은 경고 15m · 위험 8m, "
          "사람끼리 5m · 3m.", {"size": 10})],
        [("그 밖의 폰은 신호 세기로 가늠한다. 창고는 랙과 적재물 때문에 신호 세기가 거리와 "
          "늘 맞지는 않는다. 미터 대신 등급으로 알린다.", {"size": 10, "color": C.INK2})],
        [("장비를 자동으로 멈추지는 못한다. 사람에게 알리는 데까지가 이 앱의 역할이다.",
          {"size": 10, "color": C.RED})]],
       margin=0.20, space=1.3, gap=5, top=C.TOPCHIP)

r = C.body_of(s, "Result")
C.fill(r, [[("", {})]], margin=0.02)
RX, RY, RW = 6.30, 2.38, 6.27
C.text(s, RX + 0.24, RY + 0.50, RW - 0.48, 0.24,
       [[("서로 감지되는 경우의 수 — 사람이 늘수록",
          {"bold": True, "size": 10.5, "color": C.NAVY})]])
PEOPLE = [3, 5, 10, 20]
C.chart2(s, RX + 0.14, RY + 0.78, RW - 0.28, 1.92,
         [f"{p}명" for p in PEOPLE],
         [("장비를 산다면 (태그 착용자만 · 1인당 11만원)", [4 * p for p in PEOPLE]),
          ("앱을 깐다면 (모든 폰끼리 · 추가 0원)", [(4 + p) * (3 + p) // 2 for p in PEOPLE])],
         label_size=8.5, cat_size=8.5)
C.text(s, RX + 0.24, RY + 2.80, RW - 0.48, 0.22,
       [[("지게차 4대 기준. 장비 방식은 태그를 산 사람만 감지된다. 앱은 폰을 가진 모두가 "
          "서로 감지된다.", {"size": 8.5, "color": C.INK2})]])
C.text(s, RX + 0.24, RY + 3.10, RW - 0.48, 0.24,
       [[("하루에 몇 번 서로 가까워지는가", {"bold": True, "size": 10.5, "color": C.NAVY}),
         ("   WF11 서버 기록 자동 집계", {"size": 8.5, "color": C.INK3})]])
for i, (lab, desc, val, col) in enumerate([
        ("경고 등급", "경고가 뜰 만큼 가까워짐", f"하루 {M_WARN}회", C.AMBER),
        ("위험 등급", "위험이 뜰 만큼 더 가까워짐", f"하루 {M_DANGER}회", C.RED)]):
    C.kpi(s, RX + 0.20 + i * 2.98, RY + 3.34, 2.82, 0.88, "", lab, desc, val, col, est=False)
C.text(s, RX + 0.24, RY + 4.30, RW - 0.48, 0.22,
       [[(f"{M_FROM[5:]}~{M_TO[5:]} 기록이 남은 {M_DAYS}일 평균. 1건 = 가까워진 1분, 사고 건수 아님.",
          {"size": 8, "color": C.INK3})]])

# ── 4. Summary 자유 양식 — 화면과 판정 반경 ─────────────────
s = S[3]
C.fill(C.at(s, 0.78, 2.37), [[("", {})]], margin=0.02)
C.text(s, 1.00, 2.56, 11.4, 0.26,
       [[("어떻게 동작하는가", {"bold": True, "size": 13, "color": C.NAVY}),
         ("     경보 거리 · 감지되는 경우의 수 · 안전 구역", {"size": 10.5, "color": C.INK2})]])
pic(s, "fig_radius.png", 1.05, 2.98, 3.80, 648 / 358,
    "정밀 측위 폰끼리 — 지게차가 끼면 경고 15m · 위험 8m")
pic(s, "fig_pairs.png", 5.30, 3.52, 3.55, 728 / 278,
    "지게차 4대 · 사람 3명이면 21가지, 10명이면 91가지")
pic(s, "fig_zone.png", 9.30, 3.52, 3.05, 668 / 288,
    "휴게실 · 충전 구역 — 서로 울리지 않는다")
C.lines(s, 1.05, 5.62, 11.3, [
    [("·  ", {"color": C.INK3}), ("경보는 소리 · 진동 · 화면 세 가지로 동시에", {"bold": True}),
     ("  —  소음으로 못 듣거나 장갑으로 진동을 놓쳐도 나머지가 전달된다", {"color": C.INK2})],
    [("·  ", {"color": C.INK3}), ("상대가 지게차인지 사람인지, 후진 중인지까지 신호에 실어 보낸다", {"bold": True}),
     ("  —  지게차가 낀 짝은 더 먼 거리에서 울린다", {"color": C.INK2})],
    [("·  ", {"color": C.INK3}), ("같은 상대와 같은 등급이 5초 넘으면 소리를 끈다", {"bold": True}),
     ("  —  서 있을 때 계속 울리지 않는다. 더 가까워지면 즉시 다시 울린다", {"color": C.INK2})],
    [("·  ", {"color": C.INK3}), ("화면이 꺼져도 감지는 이어진다. 감지가 멈추면 알림 표시줄에 ‘이상’ 이 뜬다",
                                 {"bold": True}),
     ("  —  앱이 꺼진 줄 모르고 일하는 상황을 막는다", {"color": C.INK2})],
], size=10, gap=0.31)

# ── 5. 02. Problem Statement ────────────────────────────────
s = S[4]
C.fill(C.body_of(s, "문제의 배경"),
       [[("AS-IS", {"bold": True, "color": C.NAVY}), (" :  ", {"color": C.INK3}),
         ("우리 센터에는 지게차 · 보행자 접근 경보 장치가 없다.", {})],
        [("Pain Point", {"bold": True, "color": C.NAVY}), (" :  ", {"color": C.INK3}),
         ("지게차와 사람이 마주치는 교차 구간. 파렛트가 시야를 가린다. 소음에 경적이 묻힌다.", {})],
        [("위험했던 순간이 하루 몇 번인지 기록된 적이 없었다.", {"color": C.INK2})]],
       size=10.5, margin=0.20, space=1.35, gap=5, top=C.TOPCHIP)
C.fill(C.body_of(s, "표준/목표"),
       [[("Target", {"bold": True, "color": C.NAVY}), (" :  ", {"color": C.INK3}),
         ("서로 다가오는 것을 부딪히기 전에 알린다.", {})],
        [("Vision", {"bold": True, "color": C.NAVY}), (" :  ", {"color": C.INK3}),
         ("사람과 장비 모두가 서로의 접근을 안다.", {})]],
       size=10, margin=0.18, space=1.3, top=C.TOPCHIP)
C.fill(C.body_of(s, "현재 상황"),
       [[("Data", {"bold": True, "color": C.NAVY}), (" :  ", {"color": C.INK3}),
         ("감지 장비 ", {}), ("0대", {"bold": True, "color": C.RED}),
         (". 지게차 4대와 상시 인원 3명 모두 감지되지 않는다.", {})],
        [("Issue", {"bold": True, "color": C.NAVY}), (" :  ", {"color": C.INK3}),
         ("사고도 아차사고도 집계된 적이 없다.", {})]],
       size=10, margin=0.18, space=1.3, top=C.TOPCHIP)

why = C.body_of(s, "문제를 해결해야 하는 이유")
C.fill(why, [[("", {})]], margin=0.02)
WX, WY, WW = 0.77, 4.88, 11.80
C.text(s, WX + 0.28, WY + 0.44, WW - 0.56, 0.30,
       [[("1. 시야가 막힌다", {"bold": True, "size": 12.5, "color": C.NAVY}),
         ("      2. 경적이 안 들린다", {"bold": True, "size": 12.5, "color": C.NAVY}),
         ("      3. 주의력에만 기댄다", {"bold": True, "size": 12.5, "color": C.NAVY}),
         ("      4. 위험이 기록되지 않는다", {"bold": True, "size": 12.5, "color": C.NAVY})]],
       align=PP_ALIGN.CENTER)
C.text(s, WX + 0.28, WY + 0.88, WW - 0.56, 0.28,
       [[("앱을 깐 뒤 하루 평균 경고 등급 ", {"size": 11.5}),
         (f"{M_WARN}회", {"bold": True, "size": 11.5, "color": C.RED}),
         (" · 위험 등급 ", {"size": 11.5}),
         (f"{M_DANGER}회", {"bold": True, "size": 11.5, "color": C.RED}),
         ("가 기록됐다", {"size": 11.5})]],
       align=PP_ALIGN.CENTER)
C.text(s, WX + 0.28, WY + 1.18, WW - 0.56, 0.26,
       [[(f"WF11 {M_FROM}~{M_TO[5:]} 중 기록이 남은 {M_DAYS}일. "
          "1건 = 상대와 가까워진 1분. 사고 건수가 아니다.",
          {"size": 9.5, "color": C.INK2})]], align=PP_ALIGN.CENTER)
C.text(s, WX + 0.28, WY + 1.50, WW - 0.56, 0.28,
       [[("“지금 이 구간을 지키는 것은 장치가 아니라 ", {"italic": True, "size": 10.5, "color": C.INK2}),
         ("작업자의 주의력", {"italic": True, "bold": True, "size": 10.5}),
         ("뿐입니다.”", {"italic": True, "size": 10.5, "color": C.INK2})]],
       align=PP_ALIGN.CENTER)

# ── 6. 02. Problem Statement _ 5Why ─────────────────────────
s = S[5]
C.fill(C.named(s, "TextBox 17"),
       [[("Most likely Causes :  ", {"size": 10.5, "color": C.INK2}),
         ("지게차와 보행자가 서로 다가오는 것을 알려주는 장치가 설치돼 있지 않다",
          {"bold": True, "size": 10.5})]], margin=0.0)
C.named(s, "TextBox 17").width = Inches(11.0)
C.fill(C.named(s, "TextBox 18"),
       [[("Root Causes :  ", {"size": 10.5, "color": C.INK2}),
         ("접근 감지를 ‘장비를 사는 일’ 로만 봐 왔다. 예산이 서지 않으면 대책이 없는 상태가 이어졌다",
          {"bold": True, "size": 10.5, "color": C.RED}),
         ("   (시중 제품이 모두 태그 방식이라 — 8장)", {"size": 9, "color": C.INK3})]], margin=0.0)
C.named(s, "TextBox 18").width = Inches(11.0)

t = [sh.table for sh in s.shapes if sh.has_table][0]
WHY = [
    ("왜 교차 구간에서 위험한가?", "서로를 보지 못한 채 가까워지기 때문이다."),
    ("왜 보지 못하는가?", "적재 파렛트가 시야를 가리고, 장비 소음에 경적이 묻히기 때문이다."),
    ("왜 눈과 경적에만 의존하는가?", "접근을 대신 알려주는 장치가 설치돼 있지 않기 때문이다."),
    ("왜 설치되지 않았는가?",
     "시중 장비는 차량에 수신기, 사람에 태그를 달아야 한다. 사람이 늘 때마다 태그를 더 산다."),
    ("왜 그 방법뿐이라고 보았는가?",
     "접근 감지를 ‘장비를 사는 일’ 로만 봤다. 이미 나눠준 업무용 스마트폰이 "
     "같은 일을 할 수 있는지는 검토된 적이 없다."),
]
for i, (q, a) in enumerate(WHY):
    # 양식 1열 번호가 '1 Why ?' 와 '3. Why ?' 로 섞여 있다 — 서식은 두고 글자만 통일
    for r in t.cell(i * 2, 0).text_frame.paragraphs[0].runs[:1]:
        r.text = f"{i + 1} Why ?"
    for r in t.cell(i * 2, 0).text_frame.paragraphs[0].runs[1:]:
        r.text = ""
    C.cell(t.cell(i * 2, 1), [[(q, {"bold": True, "color": C.NAVY})]], 11, PP_ALIGN.LEFT)
    C.cell(t.cell(i * 2 + 1, 1),
           [[("→  ", {"color": C.INK3}), (a, {})]], 10.5, PP_ALIGN.LEFT)

# ── 7. 03. Tradeoffs ────────────────────────────────────────
# 양식이 3쌍짜리다. 왼쪽은 무엇을 포기했는지, 오른쪽은 포기한 쪽을 무엇으로 메웠는지.
s = S[6]
TRADE = [
    (2.37, [[("넓게 울리기  vs.  좁게 울리기", {"bold": True, "size": 12})],
            [("넓게 울리면 작업자가 앱을 꺼 둔다. 꺼진 폰은 아무것도 못 한다.",
              {"size": 10, "color": C.INK2})]],
     [[("울릴 때와 곳을 좁혔다", {"bold": True, "size": 12, "color": C.NAVY})],
      [("같은 상대 · 같은 등급이 5초 넘으면 소리를 끈다. 더 가까워지면 즉시 다시 울린다. "
        "멀어진 뒤 경고 거리에서는 다시 울리지 않는다. 휴게실 · 충전 구역에서는 울리지 않는다.",
        {"size": 9.5, "color": C.INK2})]]),
    (3.96, [[("알리기만  vs.  장비 세우기", {"bold": True, "size": 12})],
            [("자동으로 세우려면 지게차마다 정지 장치가 필요하다. 앱은 못 한다.",
              {"size": 10, "color": C.INK2})]],
     [[("알리는 것부터 시작했다", {"bold": True, "size": 12, "color": C.RED})],
      [("자동으로 멈추는 쪽이 더 확실하다. 앱을 깔기 전에는 알림도 자동 정지도 없었다. "
        "장비를 사지 않고 당장 할 수 있는 쪽을 먼저 했다. 자동 정지는 필요하면 별도 과제로 검토한다.",
        {"size": 9.5, "color": C.INK2})]]),
    (5.56, [[("미터 표시  vs.  등급 표시", {"bold": True, "size": 12})],
            [("창고는 랙 · 적재물 때문에 신호 세기가 거리와 늘 맞지는 않는다.",
              {"size": 10, "color": C.INK2})]],
     [[("경고 · 위험 등급으로 알린다", {"bold": True, "size": 12, "color": C.NAVY})],
      [("정밀 측위(UWB) 폰끼리는 실제 거리로 등급을 정한다. 그 밖의 폰은 신호 세기로 정한다. "
        "신호 세기는 미터로 바꾸지 않는다. 어느 쪽이든 작업자에게는 경고 · 위험 등급으로 알린다.",
        {"size": 9.5, "color": C.INK2})]]),
]
for y, left, right in TRADE:
    C.fill(C.at(s, 0.77, y), left, anchor=MSO_ANCHOR.MIDDLE, margin=0.22, space=1.3, gap=3)
    C.fill(C.at(s, 7.02, y), right, anchor=MSO_ANCHOR.MIDDLE, margin=0.22, space=1.3, gap=3)

# ── 8. 04. Benchmark ────────────────────────────────────────
# 우리 센터에는 감지 장비가 없다. I-PAS 는 사내 다른 조직이 쓰는 사례이므로 여기에 놓는다.
# 도입 비용도 여기서 비교한다 — 우리가 안 쓰는 것을 '절감액' 이라 부를 수는 없기 때문이다.
s = prs.slides[7]
C.fill(C.body_of(s, "대상"),
       [[("지게차 ↔ 보행자 접근 경보 장비는 시중에 이미 여러 제품이 있고, "
          "사내 다른 조직도 I-PAS 를 도입해 쓰고 있다. 우리 센터에는 없다.",
          {"bold": True, "size": 11})],
        [("감지 거리 10~25m 는 이 앱이 정밀 측위(UWB)를 쓸 때의 경고 거리 15m 와 같은 범위다.",
          {"size": 10, "color": C.INK2}),
         ("   (제품 정보 — 제조사 공개 자료 2026.09 조회, 제조사 주장 효과 수치는 인용하지 않음)",
          {"size": 8.5, "color": C.INK3})]],
       margin=0.22, anchor=MSO_ANCHOR.MIDDLE, space=1.3, top=C.TOPCHIP)
C.fill(C.body_of(s, "내용"), [[("", {})]], margin=0.02)
BX, BY, BW = 0.76, 3.48, 11.80
C.table(s, BX + 0.24, BY + 0.42, BW - 0.48,
        ["제품 · 사례", "방식", "감지 거리", "몸에 다는 것", "장비 정지", "센터 1곳 도입 비용"],
        [["ZoneSafe", "무선 태그(RFID)", "차량 사방 최대 10m", "태그 필요", "없음 (알림만)", "공개 단가 없음"],
         ["ELOKON ELOshield", "정밀 측위(UWB)", "경고 / 정지 구역 분리 · 가려져도 감지", "태그 필요",
          "자동 감속", "공개 단가 없음"],
         ["Pozyx RTLS", "정밀 측위(UWB)", "차량 사방 최대 25m", "태그 필요", "없음 (알림만)", "공개 단가 없음"],
         [[[("I-PAS  (사내 타 조직 도입)", {"bold": True})]], "태그", "3m", "태그 필요",
          "자동 정지", [[("331만원", {"bold": True, "color": C.AMBER})]]],
         [[[("SafeAlert  (이 제안)", {"bold": True, "color": C.TEAL})]], "블루투스 + 정밀 측위",
          "정밀 측위 폰끼리 15m / 8m · 그 외 신호 세기",
          [[("없음 — 지급된 폰", {"bold": True, "color": C.TEAL})]], "없음 (알림만)",
          [[("0원", {"bold": True, "color": C.TEAL})]]],
        ],
        col_w=[2.20, 1.45, 3.55, 1.40, 1.15, 1.57], row_h=0.275, head_h=0.28,
        size=9.5, head_size=9)
C.text(s, BX + 0.28, BY + 2.10, BW - 0.56, 0.24,
       [[("331만원 산출  —  ", {"bold": True, "size": 9, "color": C.NAVY}),
         ("차량 태그 52.5만 × 4대 + 자동정지장치 22만 × 4대 + 보행자 태그 11만 × 3명 (2026.08 단가). "
          "사람이 늘면 1인당 11만원 추가. 이 앱은 지급된 폰을 쓴다.",
          {"size": 9, "color": C.INK2})]])
C.text(s, BX + 0.28, BY + 2.42, BW - 0.56, 0.26,
       [[("시사점", {"bold": True, "size": 11, "color": C.NAVY})]])
C.lines(s, BX + 0.34, BY + 2.70, BW - 0.68, [
    [("·  시중 제품은 예외 없이 ", {}), ("‘차량에 수신기, 사람에 태그’", {"bold": True}),
     (" 구조다. 사람이 늘면 태그를 더 산다.", {})],
    [("·  우리 센터 규모(지게차 4대 · 상시 3명)로 장비를 갖추면 331만원, 감지되는 짝은 12가지. "
      "이 앱은 0원, 21가지.", {})],
    [("·  I-PAS 는 3m 에서 세우고, 이 앱은 15m 에서 알린다. ", {}),
     ("앱은 장비를 자동으로 세우지 못한다. 알리는 데까지가 한계다.", {"color": C.RED})],
], size=10, gap=0.26)

# ── 9. 05. Solution ─────────────────────────────────────────
s = S[8]
C.fill(C.body_of(s, "Short term"),
       [[("업무용 스마트폰에 앱을 깔아, 폰끼리 서로 감지하게 한다",
          {"bold": True, "size": 12, "color": C.NAVY})],
        [("·  정밀 측위 폰끼리는 실제 거리로 — 경고 15m · 위험 8m, 사람끼리 5m · 3m", {})],
        [("·  그 밖의 폰은 신호 세기로 판단한다 (창고에서는 거리로 환산하지 않는다)", {})],
        [("·  소리 · 진동 · 화면으로 동시에 알린다 — 하나가 막혀도 나머지가 전달된다", {})],
        [("·  상대가 지게차인지 사람인지, 후진 중인지까지 신호에 실어 보낸다", {})],
        [("·  안전 구역과 자동 소리 끄기로, 울릴 필요 없을 때는 울리지 않는다", {})],
        [("아무 감지 수단도 없던 현장에서, 지금 3개 센터가 쓰고 있다.",
          {"bold": True, "size": 10.5, "color": C.TEAL})]],
       size=10.5, margin=0.22, space=1.32, gap=4, top=C.TOPCHIP)
C.fill(C.body_of(s, "Long term"),
       [[("장비를 사지 않고, 앱 설치만으로 17개 센터로 넓힌다",
          {"bold": True, "size": 12, "color": C.NAVY})],
        [("·  9월 — WF11 · WF21 · WF25 작업자 설문으로 만족도 · 늦은 경보를 잰다", {})],
        [("·  10월 — 개인 폰 임시 설치를 끝내고 업무용 PDA 로 옮긴다", {})],
        [("·  확산 전 — 정식 과제 등록 → 보안 검토 → 안전 · 운영 승인", {})],
        [("·  승인 후 — 17개 센터로. 센터에 설치할 것 없이 앱만 깔면 된다", {})],
        [("사람이 늘어도, 센터가 늘어도 장비를 사지 않는다. 추가 비용 0원.",
          {"bold": True, "size": 10.5, "color": C.TEAL})],
        [("이 앱은 사람에게 알리는 데까지가 역할이다. 장비를 자동으로 세우는 것은 "
          "필요하면 별도 과제로 검토한다.", {"size": 10, "color": C.INK2})]],
       size=10.5, margin=0.22, space=1.32, gap=4, top=C.TOPCHIP)
C.fill(C.at(s, 0.77, 5.96),
       [[("기대 효과  —  ", {"bold": True, "size": 12, "color": C.NAVY}),
         ("감지 수단이 0 이던 현장에, 장비 구매 없이 접근 경보가 생긴다. "
          "사람이 늘어도 비용이 늘지 않는다.", {"bold": True, "size": 12})],
        [("“100% 정확한 감지가 아니라, 피할 시간을 만드는 것이 목표입니다.”",
          {"italic": True, "size": 10, "color": C.INK2})]],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, margin=0.20, space=1.3, gap=3)

# ── 10. 05. Solution — 실행 계획 ────────────────────────────
s = S[9]
t = [sh.table for sh in s.shapes if sh.has_table][0]
PLAN = [
    ("WF11 첫 배포 · 현장 실사용 시작", "제안자 · 현장", "2026.06", "완료"),
    ("작업자 의견 3차례 반영 (역할 전환 · 안전 구역 등)", "제안자 · 현장", "수시", "완료"),
    ("경보 기록 서버 저장 · 매주 자동 집계", "제안자", "2026.08", "완료"),
    ("WF11 · WF21 · WF25 세 센터 실사용", "제안자 · 현장", "진행 중", "진행"),
    ("세 센터 작업자 설문 — 만족도 · 경보가 늦었던 경험", "제안자 · 현장", "9월", "진행"),
    ("업무용 PDA 로 이전 — 개인 폰 임시 설치 종료", "제안자 · IT", "10월", "예정"),
    ("정식 과제 등록 — 확산의 첫 절차", "제안자", "확산 전", "예정"),
    ("보안 검토 — 확산 전 필수", "IT · 보안", "확산 전", "협조"),
    ("현장 사용 승인 · 확대 승인", "안전 · 운영", "확산 전", "협조"),
    ("17개 센터 확산 — 장비 구매 없이 앱 설치만", "운영 · 현장", "승인 후", "협조"),
]
STCOL = {"완료": C.TEAL, "진행": C.TEAL, "예정": C.INK2, "협조": C.AMBER}
for i, (what, who, when, st) in enumerate(PLAN, start=1):
    C.cell(t.cell(i, 0), [[(what, {})]], 10.5, PP_ALIGN.LEFT, margin=0.14)
    C.cell(t.cell(i, 1), [[(who, {})]], 10.5, PP_ALIGN.CENTER)
    C.cell(t.cell(i, 2), [[(when, {})]], 10.5, PP_ALIGN.CENTER)
    C.cell(t.cell(i, 3), [[(st, {"bold": True, "color": STCOL[st]})]], 10.5, PP_ALIGN.CENTER)

# ── 11. 05. Solution 자유 양식 — 로드맵 ─────────────────────
s = S[10]
C.fill(C.at(s, 0.78, 2.37), [[("", {})]], margin=0.02)
C.text(s, 1.00, 2.56, 11.4, 0.26,
       [[("앱을 깔기 전과 뒤 — 현장의 변화", {"bold": True, "size": 13, "color": C.NAVY}),
         ("     WF11 기준 · 장비를 사지 않고 달라진 것", {"size": 10.5, "color": C.INK2})]])
C.table(s, 1.00, 2.98, 11.35,
        ["항목", "앱 깔기 전", "앱 깐 뒤", "근거"],
        [["감지", "접근을 알려주는 장치 0대",
          [[(f"폰 {M_DEVICES}대가 서로 감지 · 소리 · 진동 · 화면 경보", {"bold": True, "color": C.TEAL})]],
          "서버 기록"],
         ["기록", "사고도 위험 순간도 기록 없음",
          [[(f"하루 평균 경고 {M_WARN}회 · 위험 {M_DANGER}회 기록", {"bold": True, "color": C.TEAL})]],
          f"기록 {M_DAYS}일"],
         ["비용", "감지 장치 없음 · 장비로 갖추려면 331만원",
          [[("추가 비용 0원 · 사람이 늘어도 0원", {"bold": True, "color": C.TEAL})]], "사내 단가"],
         ["감지 범위", "감지 장치 없음 · 태그를 사도 3명이면 12쌍",
          [[("폰 가진 모두가 서로 감지 · 3명이면 21쌍", {"bold": True, "color": C.TEAL})]], "지게차 4대"],
         ["한계", "자동 정지 장치 없음",
          [[("그대로 없음 — 사람에게 알리는 데까지", {"color": C.RED})]], "별도 과제"]],
        col_w=[1.00, 4.10, 4.95, 1.30], row_h=0.36, head_h=0.32, size=10, head_size=10,
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER])
C.text(s, 1.05, 5.20, 11.3, 0.26,
       [[("왜 되는가", {"bold": True, "size": 11, "color": C.NAVY})]])
C.lines(s, 1.10, 5.50, 11.2, [
    [("①  폰이 이미 있다", {"bold": True}),
     ("  —  살 것도, 몸에 달 것도, 센터에 설치할 것도 없다.", {"color": C.INK2})],
    [("②  사람이 늘어도 0원", {"bold": True}),
     ("  —  태그를 더 사는 대신 앱만 깐다.", {"color": C.INK2})],
    [("③  현장 의견이 바로 반영된다", {"bold": True}),
     ("  —  ‘서 있는데 계속 울린다’ → 같은 경보가 5초 넘으면 소리를 끈다.", {"color": C.INK2})],
    [("현재  ", {"bold": True, "color": C.NAVY}),
     ("WF11 · WF21 · WF25 세 센터에서 실사용 중. 실측은 WF11 뿐. 10월 업무용 PDA 로 옮긴다. "
      "보안 검토와 승인을 거쳐 17개 센터로 늘린다.", {})],
], size=10, gap=0.32)

# ── 12. 06. Metrics ─────────────────────────────────────────
# 「미발령」은 뺐다 — 정의상 관측되지 않아 지표로 성립하지 않는다.
s = S[11]
C.fill(C.body_of(s, "성공 지표"), [[("", {})]], margin=0.02, top=C.TOPCHIP)
MX, MY, MW = 0.77, 2.37, 5.84
C.text(s, MX + 0.24, MY + 0.44, MW - 0.48, 0.24,
       [[("성공 지표 네 가지", {"bold": True, "size": 10.5, "color": C.NAVY}),
         ("   밑줄 = 아직 추정", {"size": 8.5, "color": C.INK3, "u": True})]])
KPI = [
    ("지표 1", "경고 등급 경보", "경고 거리 안 (정밀 측위 · 지게차 15m)", f"하루 {M_WARN}회", C.AMBER, False),
    ("지표 2", "위험 등급 경보", "위험 거리 안 (정밀 측위 · 지게차 8m)", f"하루 {M_DANGER}회", C.RED, False),
    ("지표 3", "쓸데없이 울린 횟수", "서 있는데 울림 · 멀어졌는데 또 울림",
     "하루 3회 미만", C.INK2, True),
    ("지표 4", "작업자 만족도", "울리는 시점과 거리가 맞는가", "10점 중 7.0점", C.TEAL, True),
]
for i, (tag, name, desc, val, col, est) in enumerate(KPI):
    C.kpi(s, MX + 0.20, MY + 0.72 + i * 0.86, MW - 0.40, 0.80, tag, name, desc, val, col,
          est=est)
C.text(s, MX + 0.24, MY + 4.14, MW - 0.48, 0.44,
       [[("‘울렸어야 하는데 안 울린 횟수’ 는 뺐다", {"bold": True, "size": 9, "color": C.RED}),
         ("  —  겉으로는 아무 일도 없어 셀 수 없다. "
          "설문의 ‘경보가 늦었다’ 응답으로 대신 본다.",
          {"size": 9, "color": C.INK2})]], space=1.25)

C.fill(C.body_of(s, "지표 확보 방법"), [[("", {})]], margin=0.02, top=C.TOPCHIP)
NX = 6.74
C.text(s, NX + 0.24, MY + 0.44, MW - 0.48, 0.24,
       [[("언제 울리는가 — 시간대별 경보 수",
          {"bold": True, "size": 10.5, "color": C.NAVY}),
         (f"   WF11 · 28일 중 {M_DAYS}일 기록 · 폰 {M_DEVICES}대", {"size": 8.5, "color": C.INK3})]])
C.table(s, NX + 0.20, MY + 0.78, MW - 0.40,
        ["시간대", "경보 수", "비고"],
        [["07~09시", "157건", "—"],
         ["10시", [[("230건", {"bold": True, "color": C.RED})]],
          [[("전체의 약 30% · 가장 몰린다", {"bold": True})]]],
         ["11~12시", "174건", "—"],
         ["13~14시", "2건", "점심 · 거의 0"],
         ["15~19시", "215건", "—"],
         ["합계", [[(f"{M_WARN_T + M_DANGER_T}건", {"bold": True})]],
          f"경고 {M_WARN_T} · 위험 {M_DANGER_T}"]],
        col_w=[1.20, 1.10, 3.14], row_h=0.27, head_h=0.27, size=8.5, head_size=8.5,
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
C.lines(s, NX + 0.24, MY + 2.80, MW - 0.48, [
    [("·  경보가 울릴 때마다 시각 · 상대 · 신호 세기 · 등급이 서버에 남는다. 매주 자동 집계한다.", {})],
    [("·  같은 상대는 1분에 한 번만 기록한다. 1건 = 가까워진 1분. "
      "사고 건수가 아니라 위험했던 순간의 수다.", {})],
    [("·  ", {}), ("오전 10시에 전체의 약 30%가 몰린다. 점심은 거의 0.", {"bold": True, "color": C.NAVY}),
     (" 주의를 집중할 시간대가 처음 드러났다.", {})],
    [("·  쓸데없이 울린 횟수는 작업자 신고와 서버 기록을 대조해 잰다. "
      "만족도는 9월 세 센터 작업자 설문으로 잰다.", {})],
    [("·  사고 건수는 회사 데이터로 존재하지 않아 지표에 넣지 않았다.", {"color": C.RED})],
], size=9, gap=0.34)

# ── 13. 07. Andon ───────────────────────────────────────────
s = S[12]
C.fill(C.body_of(s, "중단 지표"), [[("", {})]], margin=0.02, top=C.TOPCHIP)
AX = 0.77
C.text(s, AX + 0.24, MY + 0.42, MW - 0.48, 0.24,
       [[("현장 운영 중", {"bold": True, "size": 10.5, "color": C.NAVY})]])
STOP, WATCH = "즉시 중단", "관측 시 판단"
FIELD = [("경보가 안 울렸다", STOP, "가까워졌는데 안 울렸다는 작업자 신고 1건이라도"),
         ("쓸데없이 계속 울린다", STOP, "서 있는데 · 멀어졌는데 울린다는 지적이 반복"),
         ("감지가 멈췄다", WATCH, "알림 표시줄에 ‘이상’ 이 뜨거나 알림이 사라짐"),
         ("작업자가 꺼 둔다", WATCH, "소리를 끄거나 앱을 꺼 둔 폰이 늘어남")]
UPD = [("옛 버전과 안 통한다", STOP, "옛 버전 폰이 새 버전 폰을 감지하지 못함"),
       ("고친 문제가 또 나온다", STOP, "이미 고친 증상을 작업자가 다시 지적함"),
       ("업데이트 중 빈틈", WATCH, "작업 시간 중 여러 폰이 한꺼번에 꺼져 있음"),
       ("업데이트 후 안 켜진다", WATCH, "설치나 역할 변경 뒤 앱이 다시 켜지지 않음")]
y = MY + 0.72
for name, tag, desc in FIELD:
    C.text(s, AX + 0.32, y, MW - 0.60, 0.22,
           [[(name, {"bold": True, "size": 10}),
             ("   " + tag, {"bold": True, "size": 8.5, "color": C.RED if tag == STOP else C.INK2})]])
    C.text(s, AX + 0.32, y + 0.19, MW - 0.60, 0.20, [[(desc, {"size": 8.5, "color": C.INK2})]])
    y += 0.42
C.text(s, AX + 0.24, y + 0.10, MW - 0.48, 0.24,
       [[("업데이트 · 수정 중", {"bold": True, "size": 10.5, "color": C.NAVY})]])
y += 0.40
for name, tag, desc in UPD:
    C.text(s, AX + 0.32, y, MW - 0.60, 0.22,
           [[(name, {"bold": True, "size": 10}),
             ("   " + tag, {"bold": True, "size": 8.5, "color": C.RED if tag == STOP else C.INK2})]])
    C.text(s, AX + 0.32, y + 0.19, MW - 0.60, 0.20, [[(desc, {"size": 8.5, "color": C.INK2})]])
    y += 0.42

C.fill(C.body_of(s, "지표 확보 방법"), [[("", {})]], margin=0.02, top=C.TOPCHIP)
C.text(s, NX + 0.24, MY + 0.42, MW - 0.48, 0.24,
       [[("무엇으로 관측하는가", {"bold": True, "size": 10.5, "color": C.NAVY})]])
C.lines(s, NX + 0.30, MY + 0.72, MW - 0.56, [
    [("·  안 울림 · 쓸데없이 울림 — 작업자 신고를 서버 경보 기록과 대조한다.", {})],
    [("·  감지 멈춤 · 안 켜짐 — 알림 표시줄을 본다. 감지가 멈추면 ‘이상’ 이 뜬다.", {})],
    [("·  옛 버전 섞임 — 배포 뒤 폰마다 버전을 확인한다. 옛 버전과도 서로 통한다.", {})],
    [("·  고친 문제가 또 나옴 — 배포 때마다 고쳤던 증상을 확인한다. 작업자 지적으로도 본다.", {})],
    [("·  꺼 두는 사람 — 작업자 설문으로 묻는다. 현장에서 폰의 소리 설정도 직접 본다.", {})],
], size=9, gap=0.325)
C.text(s, NX + 0.24, MY + 2.46, MW - 0.48, 0.24,
       [[("문제가 생기면 이렇게 되돌린다", {"bold": True, "size": 10.5, "color": C.NAVY})]])
for i, (k, v) in enumerate([
    ("즉시 되돌리기", "문제가 확인되면 바로 이전 버전으로 되돌린다"),
    ("먼저 1대에서", "모든 폰에 깔기 전에 1대에서 먼저 확인한다"),
    ("교대 때만", "작업 시간 중에는 바꾸지 않는다. 교대 전환 때만 배포한다"),
    ("다시 배포할 때", "한 번에 하나만 바꾼다. 같은 문제가 없는지 확인한 뒤 배포한다"),
    ("중단 비용 0", "앱을 꺼도 도입 전 상태로 돌아갈 뿐이다 (감지 장치 0대)"),
]):
    C.text(s, NX + 0.30, MY + 2.78 + i * 0.36, MW - 0.56, 0.34,
           [[(k, {"bold": True, "size": 9.5, "color": C.NAVY})],
            [(v, {"size": 8.5, "color": C.INK2})]], space=1.2)
# 인용 한 줄을 더 넣을 자리가 없다 — Rollback 다섯 항목이 상자를 꽉 채운다

# ── 14. 08. Feedback Loop ───────────────────────────────────
s = S[13]
C.fill(C.body_of(s, "모니터링 방법"), [[("", {})]], margin=0.02, top=C.TOPCHIP)
FX, FY, FW = 0.77, 2.38, 11.80
C.text(s, FX + 0.28, FY + 0.42, FW - 0.56, 0.24,
       [[("현장 지적 세 가지가 그대로 기능이 됐다",
          {"bold": True, "size": 11.5, "color": C.NAVY})]])
C.table(s, FX + 0.28, FY + 0.74, 5.55,
        ["차수", "반영 내용", "결과"],
        [["1차", "역할 전환 버튼 · EPJ 역할 숨김", "재설치 없이 화면에서 역할 전환"],
         ["2차", "같은 등급이 5초 넘으면 소리 끄기", "더 가까워지면 즉시 다시 울림"],
         ["3차", "휴게실 · 충전 구역을 안전 구역으로", "그 안에서는 서로 울리지 않음"]],
        col_w=[0.58, 2.87, 2.10], row_h=0.32, head_h=0.27, size=8.5, head_size=8.5,
        aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
C.table(s, FX + 6.20, FY + 0.74, 5.60,
        ["항목", "Before", "After"],
        [["앱 상태", "꺼져도 알 수 없음", "알림 표시줄에 ‘이상’ 표시"],
         ["멀어진 뒤", "멀어졌는데 다시 울림", "경고 거리에서는 다시 울리지 않음"],
         ["알림 끄기", "단계가 번거로움", "알림 본문을 누르면 즉시 무음"]],
        col_w=[1.05, 2.20, 2.35], row_h=0.32, head_h=0.27, size=8.5, head_size=8.5)
C.text(s, FX + 0.28, FY + 2.06, FW - 0.56, 0.24,
       [[("“서 있는데 계속 울린다”  “멀어졌는데 다시 울린다”  “앱이 꺼진 줄 몰랐다”",
          {"bold": True, "size": 9, "color": C.NAVY}),
         ("  —  모두 위 표의 항목으로 해결됐다. 다음 의견은 9월 설문으로 받는다.", {"size": 9})]])
C.fill(C.body_of(s, "확산 계획"),
       [[("① 9월 작업자 설문", {"bold": True, "color": C.NAVY}),
         ("  WF11 · WF21 · WF25 작업자에게 묻는다 — 경보가 늦은 적이 있는가, 만족하는가.", {})],
        [("② 업무용 PDA 이전", {"bold": True, "color": C.NAVY}),
         ("  10월, 개인 폰에 임시로 깐 것을 회사 기기로 옮긴다.", {})],
        [("③ 17개 센터 확산", {"bold": True, "color": C.NAVY}),
         ("  정식 과제 등록 → 안전 · 운영 승인. 확산 전 보안 검토를 받는다. "
          "장비 없이 앱 설치만으로 늘린다.", {})]],
       size=10.5, margin=0.24, anchor=MSO_ANCHOR.MIDDLE, space=1.35, gap=5, top=C.TOPCHIP)

prs.save(OUT)
print("wrote", OUT, f"({len(prs.slides)}장)")
